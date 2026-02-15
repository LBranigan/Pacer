#!/usr/bin/env python3
"""Generate a single hero slide PPTX for Pacer pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_DEEP = RGBColor(0x06, 0x0C, 0x1A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
DARK_CARD = RGBColor(0x16, 0x20, 0x32)
ACCENT1 = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x71, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GRAY = RGBColor(0x6B, 0x72, 0x80)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE,
             bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height):
    """Returns a textbox and its text_frame for manual rich text building."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


# ============================================================
# SLIDE - HERO
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s, DARK_BG)

# ---- RIGHT PANEL BACKGROUND ----
right_panel = add_shape(s, MSO_SHAPE.RECTANGLE,
                        Inches(6.5), Inches(0), Inches(6.833), Inches(7.5),
                        fill_color=RGBColor(0x0B, 0x14, 0x25),
                        line_color=RGBColor(0x1E, 0x29, 0x3B), line_width=0.5)

# Subtle vertical divider line
add_shape(s, MSO_SHAPE.RECTANGLE,
          Inches(6.5), Inches(0), Inches(0.015), Inches(7.5),
          fill_color=RGBColor(0x1E, 0x29, 0x3B))

# ========================================
# LEFT SIDE - Brand + Value Prop
# ========================================

# Waveform bars (logo icon) - built from rectangles
bar_x = Inches(0.8)
bar_y_center = Inches(1.55)
bar_w = Inches(0.08)
bars = [
    (0.00, 0.22, 0.5),
    (0.12, 0.38, 0.7),
    (0.24, 0.56, 1.0),
    (0.36, 0.34, 0.8),
    (0.48, 0.24, 0.6),
]
for dx, h, opacity in bars:
    bar = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                    bar_x + Inches(dx), bar_y_center - Inches(h/2),
                    bar_w, Inches(h),
                    fill_color=ACCENT1)
    bar.adjustments[0] = 0.5  # fully rounded

# Arrow head for waveform
arrow = add_shape(s, MSO_SHAPE.CHEVRON,
                  bar_x + Inches(0.58), bar_y_center - Inches(0.12),
                  Inches(0.16), Inches(0.24),
                  fill_color=ACCENT2)

# "PACER" wordmark
add_text(s, "PACER", Inches(1.65), Inches(1.3), Inches(2.5), Inches(0.5),
         font_size=26, color=ACCENT2, bold=True, font_name='Calibri')

# Main headline - multi-line with gradient effect
# Line 1
add_text(s, "AI that hears", Inches(0.8), Inches(2.2), Inches(5.2), Inches(0.7),
         font_size=46, color=WHITE, bold=True, font_name='Calibri')
# Line 2 (gradient accent)
add_text(s, "how students", Inches(0.8), Inches(2.9), Inches(5.2), Inches(0.7),
         font_size=46, color=ACCENT2, bold=True, font_name='Calibri')
# Line 3 (gradient accent continued)
add_text(s, "struggle to read", Inches(0.8), Inches(3.6), Inches(5.2), Inches(0.7),
         font_size=46, color=ACCENT2, bold=True, font_name='Calibri')

# Description paragraph
add_text(s, "Students read aloud from real books. Pacer listens, detects every moment of struggle, and returns rich diagnostic data to teachers \u2014 with zero prep.",
         Inches(0.8), Inches(4.6), Inches(5.0), Inches(1.0),
         font_size=15, color=MUTED, font_name='Calibri')

# Feature pills
pills = [
    ("Middle School RTI", ACCENT1),
    ("13 Miscue Types", GREEN),
    ("3-Engine ASR", PURPLE),
    ("Real Books, Not Screens", ORANGE),
]
pill_x = Inches(0.8)
pill_y = Inches(5.85)
for i, (label, dot_color) in enumerate(pills):
    pw = Inches(2.4) if i == 3 else Inches(2.0)
    # Pill background
    pill_bg = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                        pill_x, pill_y, pw, Inches(0.38),
                        fill_color=DARK_CARD,
                        line_color=RGBColor(0x33, 0x41, 0x55), line_width=0.75)
    pill_bg.adjustments[0] = 0.5
    # Dot
    add_shape(s, MSO_SHAPE.OVAL,
              pill_x + Inches(0.15), pill_y + Inches(0.13),
              Inches(0.12), Inches(0.12),
              fill_color=dot_color)
    # Label
    add_text(s, label, pill_x + Inches(0.35), pill_y + Inches(0.02),
             pw - Inches(0.45), Inches(0.34),
             font_size=10, color=MUTED, bold=True, font_name='Calibri')
    pill_x += pw + Inches(0.12)
    if i == 1:  # wrap to second row
        pill_x = Inches(0.8)
        pill_y += Inches(0.48)


# ========================================
# RIGHT SIDE - Demo Panel (Assessment Mockup)
# ========================================

# Demo window container
demo_x = Inches(7.0)
demo_y = Inches(0.7)
demo_w = Inches(5.8)
demo_h = Inches(6.1)

# Window frame
demo_frame = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                       demo_x, demo_y, demo_w, demo_h,
                       fill_color=DARK_DEEP,
                       line_color=RGBColor(0x2A, 0x35, 0x4A), line_width=1)
demo_frame.adjustments[0] = 0.03
# glow effect (larger shadow shape behind)
glow = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                 demo_x - Inches(0.05), demo_y - Inches(0.05),
                 demo_w + Inches(0.1), demo_h + Inches(0.1),
                 fill_color=None,
                 line_color=RGBColor(0x08, 0x91, 0xB2), line_width=0.5)
glow.adjustments[0] = 0.03
# Move glow behind
# (in PPTX, shapes are z-ordered by creation - we need to reorder)
# We'll just leave it as a subtle border effect

# Title bar
titlebar = add_shape(s, MSO_SHAPE.RECTANGLE,
                     demo_x, demo_y, demo_w, Inches(0.45),
                     fill_color=RGBColor(0x10, 0x1A, 0x2E))
# Traffic light dots
for i, color in enumerate([RED, RGBColor(0xEA, 0xB3, 0x08), GREEN]):
    add_shape(s, MSO_SHAPE.OVAL,
              demo_x + Inches(0.2 + i * 0.22), demo_y + Inches(0.14),
              Inches(0.13), Inches(0.13),
              fill_color=color)
# Title bar text
add_text(s, "pacer \u2014 assessment results",
         demo_x + Inches(0.9), demo_y + Inches(0.07),
         Inches(3), Inches(0.35),
         font_size=9, color=DIM, font_name='Consolas')

# ---- PASSAGE LABEL ----
label_y = demo_y + Inches(0.6)
add_text(s, "STUDENT READING \u2014 PASSAGE ANALYSIS",
         demo_x + Inches(0.3), label_y, Inches(4), Inches(0.3),
         font_size=8, color=DIM, bold=True, font_name='Calibri')

# ---- COLOR-CODED WORDS ----
# We'll lay out each word individually with precise positioning
word_start_y = label_y + Inches(0.4)
word_x_start = demo_x + Inches(0.3)
max_word_x = demo_x + demo_w - Inches(0.3)

# Word data: (text, type, annotation_text, annotation_color)
# Types: correct, hesitation, struggle, substitution, pause, selfcorr, omission
words = [
    ("The", "correct", None, None),
    ("boy", "correct", None, None),
    ("walked", "hesitation", None, None),
    ("slowly", "correct", None, None),
    ("to", "correct", None, None),
    ("the", "correct", None, None),
    ("gro\u2013groc\u2013grocery", "struggle", "3.2s \u2022 decoding", TEAL),
    ("store", "correct", None, None),
    ("and", "correct", None, None),
    ("bringed", "substitution", "bought \u2192", ORANGE),
    ("some", "correct", None, None),
    ("\u22ef 4.1s", "pause", None, None),
    ("milk", "selfcorr", "self-corrected", PURPLE),
    ("for", "correct", None, None),
    ("his", "correct", None, None),
    ("grandmother", "omission", None, None),
    ("who", "correct", None, None),
    ("lived", "correct", None, None),
    ("around", "substitution", "across \u2192", ORANGE),
    ("the", "correct", None, None),
    ("street.", "correct", None, None),
]

# Color map for word types
type_colors = {
    "correct": RGBColor(0xC8, 0xD0, 0xDB),
    "hesitation": ORANGE,
    "struggle": TEAL,
    "substitution": ORANGE,
    "pause": GRAY,
    "selfcorr": PURPLE,
    "omission": RED,
}

# Approximate character widths for layout (in inches, for ~14pt font)
CHAR_W = 0.09
WORD_GAP = Inches(0.18)
LINE_H = Inches(0.55)  # line height including annotation space
ANNOTATION_OFFSET = Inches(-0.22)

cur_x = word_x_start
cur_y = word_start_y
line_num = 0

for word_text, word_type, annotation, ann_color in words:
    # Estimate word width
    w = max(len(word_text) * CHAR_W, 0.3)
    word_w = Inches(w)

    # Check line wrap
    if cur_x + word_w > max_word_x:
        cur_x = word_x_start
        cur_y += LINE_H
        line_num += 1

    # Draw annotation above the word if present
    if annotation:
        add_text(s, annotation,
                 cur_x, cur_y + ANNOTATION_OFFSET,
                 word_w + Inches(0.3), Inches(0.2),
                 font_size=7, color=ann_color or DIM, font_name='Consolas')

    # Draw the word
    word_color = type_colors.get(word_type, WHITE)
    is_bold = word_type in ("struggle", "substitution", "omission")

    if word_type == "pause":
        # Special: pause indicator with background
        pause_w = Inches(0.8)
        pause_bg = add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
                             cur_x, cur_y + Inches(0.02),
                             pause_w, Inches(0.28),
                             fill_color=RGBColor(0x1A, 0x1F, 0x2E),
                             line_color=RGBColor(0x33, 0x3D, 0x50), line_width=0.5)
        pause_bg.adjustments[0] = 0.3
        add_text(s, word_text, cur_x + Inches(0.08), cur_y + Inches(0.0),
                 pause_w - Inches(0.1), Inches(0.3),
                 font_size=10, color=GRAY, font_name='Consolas', alignment=PP_ALIGN.CENTER)
        cur_x += pause_w + WORD_GAP
        continue

    if word_type == "omission":
        # Strikethrough effect: draw word + line through it
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.1), Inches(0.32),
                 font_size=14, color=RGBColor(0xEF, 0x44, 0x44), font_name='Calibri')
        # Strikethrough line
        strike_line = add_shape(s, MSO_SHAPE.RECTANGLE,
                                cur_x + Inches(0.02), cur_y + Inches(0.16),
                                word_w - Inches(0.02), Inches(0.02),
                                fill_color=RGBColor(0xEF, 0x44, 0x44))
    elif word_type == "struggle":
        # Dotted underline effect: word + dots below
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.5), Inches(0.32),
                 font_size=14, color=TEAL, bold=True, font_name='Calibri')
        # Dotted underline (series of small circles)
        actual_w = min(word_w + Inches(0.4), Inches(1.8))
        for dot_i in range(int(actual_w / Inches(0.07))):
            if dot_i % 2 == 0:
                add_shape(s, MSO_SHAPE.OVAL,
                          cur_x + Inches(dot_i * 0.07), cur_y + Inches(0.3),
                          Inches(0.04), Inches(0.04),
                          fill_color=RGBColor(0x14, 0xB8, 0xA6))
        word_w += Inches(0.5)  # wider for this long word
    elif word_type == "substitution":
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.1), Inches(0.32),
                 font_size=14, color=ORANGE, bold=True, font_name='Calibri')
        # Solid underline
        add_shape(s, MSO_SHAPE.RECTANGLE,
                  cur_x, cur_y + Inches(0.3),
                  word_w, Inches(0.025),
                  fill_color=RGBColor(0xF9, 0x71, 0x16))
    elif word_type == "hesitation":
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.1), Inches(0.32),
                 font_size=14, color=ORANGE, font_name='Calibri')
        # Dashed underline
        dash_count = max(int(word_w / Inches(0.1)), 3)
        for d in range(dash_count):
            if d % 2 == 0:
                add_shape(s, MSO_SHAPE.RECTANGLE,
                          cur_x + Inches(d * 0.1), cur_y + Inches(0.3),
                          Inches(0.06), Inches(0.02),
                          fill_color=RGBColor(0xF9, 0x71, 0x16))
    elif word_type == "selfcorr":
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.1), Inches(0.32),
                 font_size=14, color=PURPLE, font_name='Calibri')
    else:
        # correct
        add_text(s, word_text, cur_x, cur_y,
                 word_w + Inches(0.1), Inches(0.32),
                 font_size=14, color=RGBColor(0xA0, 0xAA, 0xBA), font_name='Calibri')

    cur_x += word_w + WORD_GAP


# ---- METRICS BAR ----
metrics_y = demo_y + demo_h - Inches(1.3)

# Divider line
add_shape(s, MSO_SHAPE.RECTANGLE,
          demo_x + Inches(0.3), metrics_y,
          demo_w - Inches(0.6), Inches(0.015),
          fill_color=RGBColor(0x1E, 0x29, 0x3B))

metrics = [
    ("62", "WCPM", ACCENT2),
    ("78%", "ACCURACY", ACCENT2),
    ("4", "ERRORS", ORANGE),
    ("1", "STRUGGLE", TEAL),
    ("1", "SELF-CORR", PURPLE),
    ("1", "OMISSION", RED),
]

for i, (val, label, color) in enumerate(metrics):
    mx = demo_x + Inches(0.3 + i * 0.92)
    add_text(s, val, mx, metrics_y + Inches(0.15),
             Inches(0.85), Inches(0.4),
             font_size=22, color=color, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
    add_text(s, label, mx, metrics_y + Inches(0.55),
             Inches(0.85), Inches(0.25),
             font_size=7, color=DIM, bold=True, font_name='Calibri', alignment=PP_ALIGN.CENTER)


# ---- LEGEND BAR at bottom of demo ----
legend_y = demo_y + demo_h - Inches(0.45)
legend_items = [
    ("\u25A0 Omission", RED),
    ("\u25A0 Substitution", ORANGE),
    ("\u25A0 Struggle", TEAL),
    ("\u25A0 Self-Correction", PURPLE),
    ("\u25A0 Pause", GRAY),
    ("\u25A0 Hesitation", ORANGE),
]
for i, (label, color) in enumerate(legend_items):
    lx = demo_x + Inches(0.3 + i * 0.95)
    add_text(s, label, lx, legend_y,
             Inches(1.0), Inches(0.25),
             font_size=7, color=color, font_name='Calibri')


# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pitch-deck-hero.pptx')
prs.save(out)
print(f"Saved to {out}")
