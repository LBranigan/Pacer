#!/usr/bin/env python3
"""Generate 5 emotional hero slide concepts for Pacer pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, random

# Colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_DEEP = RGBColor(0x06, 0x0C, 0x1A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
DARK_CARD = RGBColor(0x16, 0x20, 0x32)
ACCENT1 = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
DARKER_DIM = RGBColor(0x33, 0x41, 0x55)
RED = RGBColor(0xEF, 0x44, 0x44)
SOFT_RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xF9, 0x71, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WARM_DARK = RGBColor(0x1A, 0x15, 0x10)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txt(slide, text, left, top, width, height, size=18, color=WHITE,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return txBox


def multiline(slide, lines, left, top, width, height, align=PP_ALIGN.LEFT, spacing=6):
    """lines: list of (text, size, color, bold, italic, font)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, italic, font) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return txBox


def rect(slide, left, top, width, height, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def rrect(slide, left, top, width, height, fill=None, line=None, lw=None, r=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = r
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def circle(slide, left, top, size, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def pacer_brand(slide, x, y, size='small'):
    """Add Pacer waveform + wordmark."""
    # Waveform bars
    bw = Inches(0.06) if size == 'small' else Inches(0.08)
    gap = 0.09 if size == 'small' else 0.12
    bars = [(0, 0.18), (1, 0.3), (2, 0.45), (3, 0.28), (4, 0.2)]
    max_h = 0.45
    for i, (_, h) in enumerate(bars):
        hs = h if size == 'small' else h * 1.3
        ms = max_h if size == 'small' else max_h * 1.3
        bx = x + Inches(i * gap)
        by = y + Inches((ms - hs) / 2)
        bar = rrect(slide, bx, by, bw, Inches(hs), fill=ACCENT2, r=0.5)
    fs = 20 if size == 'small' else 28
    txt(slide, "PACER", x + Inches(len(bars) * gap + 0.08), y - Inches(0.05),
        Inches(2), Inches(0.5), size=fs, color=ACCENT2, bold=True)


# ============================================================
# CONCEPT 1: "THE INVISIBLE CLASSROOM"
# 30 student dots. 12 glow red. "Their teacher doesn't know."
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, DARK_DEEP)

txt(s1, "CONCEPT 1: The Invisible Classroom", Inches(10), Inches(0.3), Inches(3), Inches(0.3),
    size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)

pacer_brand(s1, Inches(0.7), Inches(0.4))

# Main headline
txt(s1, "One classroom. Thirty students.", Inches(0.7), Inches(1.3), Inches(8), Inches(0.7),
    size=42, color=WHITE, bold=True, font='Calibri')

# The classroom grid — 6 rows x 5 cols = 30 students
# 12 of them are "struggling" but invisible
grid_x = Inches(1.0)
grid_y = Inches(2.5)
dot_size = Inches(0.55)
gap_x = Inches(0.35)
gap_y = Inches(0.2)

# Which students are struggling (12 out of 30)
struggling = {1, 3, 6, 8, 11, 14, 16, 19, 22, 24, 27, 29}

for i in range(30):
    row = i // 6
    col = i % 6
    cx = grid_x + col * (dot_size + gap_x)
    cy = grid_y + row * (dot_size + gap_y)

    if i in struggling:
        # Outer glow ring
        circle(s1, cx - Inches(0.06), cy - Inches(0.06),
               dot_size + Inches(0.12),
               fill=RGBColor(0x2A, 0x08, 0x08), line=RGBColor(0xEF, 0x44, 0x44), lw=1.5)
        # Inner dot
        circle(s1, cx, cy, dot_size,
               fill=RGBColor(0x3A, 0x10, 0x10), line=RED, lw=1)
        # Person icon (simple)
        txt(s1, "\U0001F9D1", cx + Inches(0.08), cy + Inches(0.06), dot_size, dot_size,
            size=20, color=RED, align=PP_ALIGN.CENTER)
    else:
        circle(s1, cx, cy, dot_size,
               fill=RGBColor(0x12, 0x18, 0x28), line=RGBColor(0x25, 0x30, 0x45), lw=0.75)
        txt(s1, "\U0001F9D1", cx + Inches(0.08), cy + Inches(0.06), dot_size, dot_size,
            size=20, color=DIM, align=PP_ALIGN.CENTER)

# Right side — the emotional text
right_x = Inches(6.8)

multiline(s1, [
    ("Twelve are struggling", 36, RED, True, False, 'Calibri'),
    ("to read.", 36, RED, True, False, 'Calibri'),
], right_x, Inches(2.5), Inches(6), Inches(1.2), spacing=2)

multiline(s1, [
    ("Their teacher doesn't know.", 28, MUTED, False, False, 'Calibri'),
    ("", 10, MUTED, False, False, 'Calibri'),
    ("They pass basic screenings. They memorize words.", 16, DIM, False, True, 'Calibri'),
    ("They guess from context. They devote 90% of their", 16, DIM, False, True, 'Calibri'),
    ("brainpower to decoding \u2014 and 0% to comprehension.", 16, DIM, False, True, 'Calibri'),
], right_x, Inches(3.9), Inches(6), Inches(2.5), spacing=4)

# Bottom callout
rrect(s1, right_x, Inches(6.0), Inches(5.5), Inches(0.7),
      fill=RGBColor(0x08, 0x1A, 0x28), line=ACCENT1, lw=1, r=0.08)
txt(s1, "Pacer finds them.", right_x + Inches(0.3), Inches(6.08), Inches(4), Inches(0.5),
    size=22, color=ACCENT2, bold=True)


# ============================================================
# CONCEPT 2: "EVERY WORD IS A BATTLE"
# One giant word. Radiating struggle data. The cognitive load of one word.
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, DARK_DEEP)

txt(s2, "CONCEPT 2: Every Word Is a Battle", Inches(10), Inches(0.3), Inches(3), Inches(0.3),
    size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)

pacer_brand(s2, Inches(0.7), Inches(0.4))

# Subtitle context
txt(s2, "For a struggling reader, this is what one word looks like:",
    Inches(0.7), Inches(1.3), Inches(12), Inches(0.5),
    size=18, color=MUTED, align=PP_ALIGN.CENTER)

# THE BIG WORD — center of slide
# Outer glow rings
cx, cy = Inches(6.666), Inches(3.75)
for ring_size, ring_color in [
    (Inches(4.5), RGBColor(0x12, 0x08, 0x08)),
    (Inches(3.5), RGBColor(0x1A, 0x0A, 0x0A)),
    (Inches(2.6), RGBColor(0x22, 0x0E, 0x0E)),
]:
    circle(s2, cx - ring_size/2, cy - ring_size/2, ring_size,
           fill=ring_color, line=RGBColor(0x3A, 0x15, 0x15), lw=0.5)

# Failed attempts scattered around the word
attempts = [
    ("gr\u2013", Inches(3.0), Inches(2.2), 22, RGBColor(0x80, 0x30, 0x30)),
    ("gro\u2013", Inches(4.2), Inches(2.0), 20, RGBColor(0x90, 0x35, 0x35)),
    ("groc\u2013", Inches(8.5), Inches(2.0), 20, RGBColor(0xA0, 0x3A, 0x3A)),
    ("gros\u2013 ?", Inches(9.5), Inches(2.5), 18, RGBColor(0xB0, 0x40, 0x40)),
    ("groshy?", Inches(9.8), Inches(4.2), 16, RGBColor(0xC0, 0x45, 0x45)),
    ("gro\u2013", Inches(3.2), Inches(4.8), 18, RGBColor(0x90, 0x35, 0x35)),
]
for attempt_text, ax, ay, asize, acolor in attempts:
    txt(s2, attempt_text, ax, ay, Inches(1.5), Inches(0.5),
        size=asize, color=acolor, italic=True, font='Georgia')

# The actual word, huge
txt(s2, "grocery", Inches(4.1), Inches(3.05), Inches(5.2), Inches(1.2),
    size=72, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font='Calibri')

# Radiating data labels
labels = [
    ("\u23f1 3.2 seconds of silence", Inches(1.0), Inches(3.0), AMBER),
    ("\U0001F504 two failed attempts", Inches(0.5), Inches(3.8), ORANGE),
    ("\u274C gave up and guessed", Inches(0.7), Inches(4.5), RED),
    ("\U0001F9E0 all brainpower spent here", Inches(9.2), Inches(3.5), PURPLE),
    ("\U0001F4C9 comprehension = zero", Inches(9.5), Inches(4.8), RED),
]
for label_text, lx, ly, lcolor in labels:
    txt(s2, label_text, lx, ly, Inches(3.5), Inches(0.4),
        size=14, color=lcolor, bold=True)

# Bottom text
txt(s2, "For 40% of middle schoolers, this is every word.",
    Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
    size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(s2, "Pacer hears every hesitation, every failed attempt, every silent struggle \u2014 and logs it automatically.",
    Inches(2), Inches(6.4), Inches(9.3), Inches(0.4),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# CONCEPT 3: "TWO PATHS"
# A student's life diverges based on whether they're identified.
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, DARK_DEEP)

txt(s3, "CONCEPT 3: Two Paths", Inches(10), Inches(0.3), Inches(3), Inches(0.3),
    size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)

pacer_brand(s3, Inches(0.7), Inches(0.4))

# "Same student. Same struggle. Different outcome."
txt(s3, "Same student. Same struggle.", Inches(0.7), Inches(1.2), Inches(12), Inches(0.6),
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s3, "Different outcome.", Inches(0.7), Inches(1.85), Inches(12), Inches(0.6),
    size=40, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

# Central student dot (the fork point)
fork_x = Inches(6.666)
fork_y = Inches(3.0)
circle(s3, fork_x - Inches(0.35), fork_y - Inches(0.35), Inches(0.7),
       fill=ACCENT1, line=ACCENT2, lw=2)
txt(s3, "\U0001F9D2", fork_x - Inches(0.2), fork_y - Inches(0.2), Inches(0.4), Inches(0.4),
    size=22, align=PP_ALIGN.CENTER)
txt(s3, "6th Grade", fork_x - Inches(0.6), fork_y + Inches(0.4), Inches(1.2), Inches(0.3),
    size=10, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)

# --- LEFT PATH: Unidentified (downward, red) ---
path_steps_left = [
    (Inches(2.8), Inches(3.5), "7th grade", "Still struggling. Nobody notices.", RED),
    (Inches(1.5), Inches(4.3), "8th grade", "Falls further behind. Gives up on reading.", RED),
    (Inches(0.5), Inches(5.1), "High school", "Drops out or barely graduates.", SOFT_RED),
    (Inches(0.3), Inches(5.9), "Adulthood", "Limited opportunities. The gap never closes.", RGBColor(0x80, 0x30, 0x30)),
]

# Draw connecting lines for left path
for i in range(len(path_steps_left)):
    sx, sy, _, _, color = path_steps_left[i]
    if i == 0:
        # From fork to first step
        rect(s3, fork_x - Inches(0.5), fork_y + Inches(0.3), Inches(0.04), Inches(0.5),
             fill=RGBColor(0x3A, 0x15, 0x15))
    if i < len(path_steps_left) - 1:
        nx, ny = path_steps_left[i+1][0], path_steps_left[i+1][1]
        rect(s3, sx + Inches(0.15), sy + Inches(0.3), Inches(0.04), Inches(0.5),
             fill=RGBColor(0x2A, 0x10, 0x10))

for sx, sy, grade, desc, color in path_steps_left:
    circle(s3, sx, sy, Inches(0.3), fill=RGBColor(0x2A, 0x0A, 0x0A), line=color, lw=1)
    txt(s3, grade, sx + Inches(0.4), sy - Inches(0.05), Inches(1.5), Inches(0.25),
        size=12, color=color, bold=True)
    txt(s3, desc, sx + Inches(0.4), sy + Inches(0.18), Inches(3), Inches(0.25),
        size=11, color=DIM, italic=True)

# "WITHOUT PACER" label
txt(s3, "\u2718  WITHOUT IDENTIFICATION", Inches(0.3), Inches(2.8), Inches(3), Inches(0.3),
    size=11, color=RED, bold=True)

# --- RIGHT PATH: Identified (upward, green) ---
path_steps_right = [
    (Inches(9.5), Inches(3.3), "Identified", "Pacer flags consistent decoding struggle.", ACCENT2),
    (Inches(10.0), Inches(4.0), "RTI Tier 2", "Targeted phonics intervention begins.", GREEN),
    (Inches(10.3), Inches(4.7), "8th grade", "Reading at grade level. Comprehending.", GREEN),
    (Inches(10.5), Inches(5.4), "High school", "Confident. Engaged. Succeeding.", RGBColor(0x16, 0xA3, 0x4A)),
]

# Connecting lines for right path
for i in range(len(path_steps_right)):
    sx, sy = path_steps_right[i][0], path_steps_right[i][1]
    if i == 0:
        rect(s3, fork_x + Inches(0.4), fork_y + Inches(0.1), Inches(0.04), Inches(0.5),
             fill=RGBColor(0x0A, 0x2A, 0x1A))
    if i < len(path_steps_right) - 1:
        rect(s3, sx + Inches(0.15), sy + Inches(0.3), Inches(0.04), Inches(0.4),
             fill=RGBColor(0x0A, 0x25, 0x15))

for sx, sy, grade, desc, color in path_steps_right:
    circle(s3, sx, sy, Inches(0.3), fill=RGBColor(0x0A, 0x1E, 0x12), line=color, lw=1)
    txt(s3, grade, sx + Inches(0.4), sy - Inches(0.05), Inches(2), Inches(0.25),
        size=12, color=color, bold=True)
    txt(s3, desc, sx + Inches(0.4), sy + Inches(0.18), Inches(2.5), Inches(0.25),
        size=11, color=MUTED, italic=True)

txt(s3, "\u2714  WITH PACER", Inches(9.5), Inches(2.8), Inches(3), Inches(0.3),
    size=11, color=GREEN, bold=True)

# Bottom
txt(s3, "The only difference is whether someone noticed in time.",
    Inches(2), Inches(6.6), Inches(9.3), Inches(0.5),
    size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s3, "Pacer notices.", Inches(2), Inches(7.0), Inches(9.3), Inches(0.3),
    size=18, color=ACCENT2, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# CONCEPT 4: "SHE'S NOT SLOW. SHE'S DROWNING."
# Pure typography. Massive emotional impact. Scattered letters.
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, DARK_DEEP)

txt(s4, "CONCEPT 4: Typography-Driven", Inches(10), Inches(0.3), Inches(3), Inches(0.3),
    size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)

# Scattered background letters (the cognitive chaos of a struggling reader)
random.seed(42)
scatter_chars = list("comprehensionreadingfluencystruggledecodingphonicswordslettersmeaning")
for ch in scatter_chars[:50]:
    sx = Inches(random.uniform(0, 12.5))
    sy = Inches(random.uniform(0.5, 7))
    ssize = random.randint(14, 40)
    opacity_val = random.randint(8, 20)
    scolor = RGBColor(opacity_val, opacity_val + 4, opacity_val + 12)
    txt(s4, ch, sx, sy, Inches(0.5), Inches(0.5),
        size=ssize, color=scolor, font='Georgia')

# Main text — massive, centered, raw
txt(s4, "She's not slow.", Inches(1), Inches(1.5), Inches(11.3), Inches(1),
    size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font='Calibri')

txt(s4, "She's drowning.", Inches(1), Inches(2.7), Inches(11.3), Inches(1),
    size=60, color=RED, bold=True, align=PP_ALIGN.CENTER, font='Calibri')

# The explanation — the 90/0 split
multiline(s4, [
    ("90% of her brainpower goes to decoding words.", 22, MUTED, False, False, 'Calibri'),
    ("0% is left for comprehension.", 22, MUTED, False, False, 'Calibri'),
    ("", 12, DIM, False, False, 'Calibri'),
    ("She's memorized enough words to pass basic screenings.", 16, DIM, False, True, 'Calibri'),
    ("Her teacher thinks she's \"a little slow, but getting the words right.\"", 16, DIM, False, True, 'Calibri'),
    ("", 12, DIM, False, False, 'Calibri'),
    ("She is one of 12 million American students stuck in this invisible gap.", 16, DIM, False, True, 'Calibri'),
], Inches(2.5), Inches(4.0), Inches(8.3), Inches(2.5), align=PP_ALIGN.CENTER, spacing=3)

# PACER line
rrect(s4, Inches(3.5), Inches(6.4), Inches(6.3), Inches(0.65),
      fill=RGBColor(0x08, 0x1A, 0x28), line=ACCENT1, lw=1.5, r=0.1)

pacer_brand(s4, Inches(4.4), Inches(6.48))
txt(s4, "finds her.", Inches(6.2), Inches(6.48), Inches(3), Inches(0.45),
    size=20, color=WHITE, bold=True)


# ============================================================
# CONCEPT 5: "THE BOOK AND THE SIGNAL"
# Physical book on left, glowing with hidden data.
# A warm, hopeful tone. "Inside every reading session is a signal."
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, RGBColor(0x0A, 0x0E, 0x18))

txt(s5, "CONCEPT 5: The Book & The Signal", Inches(10), Inches(0.3), Inches(3), Inches(0.3),
    size=9, color=DIM, bold=True, align=PP_ALIGN.RIGHT)

# --- LEFT HALF: A beautiful book, warm tones ---
# Warm background wash behind book
circle(s5, Inches(0.5), Inches(0.5), Inches(6.5),
       fill=RGBColor(0x14, 0x11, 0x0C))

# Book - large, prominent
bk_x, bk_y = Inches(1.2), Inches(1.0)
bk_w, bk_h = Inches(4.0), Inches(5.2)

# Book shadow
rrect(s5, bk_x + Inches(0.12), bk_y + Inches(0.12), bk_w, bk_h,
      fill=RGBColor(0x06, 0x04, 0x02), r=0.02)
# Page edges
rrect(s5, bk_x + Inches(0.15), bk_y + Inches(0.1), bk_w - Inches(0.15), bk_h - Inches(0.1),
      fill=RGBColor(0xFA, 0xF7, 0xF0), r=0.01)
# Book cover
rrect(s5, bk_x, bk_y, bk_w, bk_h,
      fill=RGBColor(0x6B, 0x4E, 0x35), line=RGBColor(0x55, 0x3D, 0x28), lw=1.5, r=0.02)
# Spine
rect(s5, bk_x, bk_y, Inches(0.25), bk_h,
     fill=RGBColor(0x55, 0x3D, 0x28))

# Title on cover
rrect(s5, bk_x + Inches(0.6), bk_y + Inches(0.7), Inches(2.8), Inches(1.0),
      fill=RGBColor(0x5A, 0x42, 0x2E), line=RGBColor(0x8A, 0x70, 0x56), lw=0.75, r=0.04)
txt(s5, "Number the Stars", bk_x + Inches(0.7), bk_y + Inches(0.75), Inches(2.6), Inches(0.45),
    size=18, color=RGBColor(0xFA, 0xF0, 0xE0), bold=True, align=PP_ALIGN.CENTER, font='Georgia')
txt(s5, "Lois Lowry", bk_x + Inches(0.7), bk_y + Inches(1.2), Inches(2.6), Inches(0.35),
    size=12, color=RGBColor(0xC4, 0xAE, 0x96), italic=True, align=PP_ALIGN.CENTER, font='Georgia')

# Open page visualization with text lines
page_x = bk_x + Inches(0.5)
page_y = bk_y + Inches(2.0)
for i in range(8):
    ly = page_y + Inches(i * 0.3)
    line_w = Inches(random.uniform(2.0, 3.0))
    rect(s5, page_x, ly, line_w, Inches(0.06),
         fill=RGBColor(0x45, 0x38, 0x2A))

# Glowing signal emanating from book (circles radiating outward)
for ring_r, ring_opacity in [(2.8, 0x14), (3.3, 0x10), (3.8, 0x0C)]:
    circle(s5, bk_x + Inches(2.0) - Inches(ring_r/2),
           bk_y + Inches(2.6) - Inches(ring_r/2),
           Inches(ring_r),
           fill=None, line=RGBColor(0x08, ring_opacity + 0x40, ring_opacity + 0x60), lw=0.5)

# "A real book in their hands" label
txt(s5, "A real book in their hands.", Inches(0.7), Inches(6.5), Inches(5), Inches(0.4),
    size=16, color=RGBColor(0x8A, 0x76, 0x60), italic=True, align=PP_ALIGN.CENTER)

# --- RIGHT HALF: The signal Pacer extracts ---
right_x = Inches(6.0)

txt(s5, "Inside every reading session", right_x + Inches(0.5), Inches(1.2), Inches(6.5), Inches(0.5),
    size=34, color=WHITE, bold=True)
txt(s5, "is a signal.", right_x + Inches(0.5), Inches(1.85), Inches(6.5), Inches(0.5),
    size=34, color=ACCENT2, bold=True)

multiline(s5, [
    ("A hesitation before a hard word.", 18, ORANGE, False, False, 'Calibri'),
    ("A substitution that reveals a decoding gap.", 18, ORANGE, False, False, 'Calibri'),
    ("A repetition that shows lost meaning.", 18, PURPLE, False, False, 'Calibri'),
    ("A silence where confidence used to be.", 18, GRAY, False, False, 'Calibri'),
], right_x + Inches(0.5), Inches(2.7), Inches(6), Inches(2.5), spacing=12)

multiline(s5, [
    ("Teachers can't hear it with 120 students.", 15, DIM, False, True, 'Calibri'),
    ("Standardized tests can't measure it.", 15, DIM, False, True, 'Calibri'),
    ("WCPM scores don't capture it.", 15, DIM, False, True, 'Calibri'),
], right_x + Inches(0.5), Inches(4.6), Inches(6), Inches(1.5), spacing=6)

# PACER hears it
rrect(s5, right_x + Inches(0.5), Inches(5.8), Inches(5.5), Inches(0.85),
      fill=RGBColor(0x08, 0x1A, 0x28), line=ACCENT1, lw=1.5, r=0.06)

pacer_brand(s5, right_x + Inches(0.8), Inches(5.92), size='big')
txt(s5, "hears it.", right_x + Inches(2.4), Inches(5.9), Inches(3), Inches(0.55),
    size=26, color=WHITE, bold=True)

txt(s5, "AI-Powered Reading Struggle Detection for Middle School",
    right_x + Inches(0.5), Inches(6.8), Inches(5.5), Inches(0.3),
    size=12, color=DIM)


# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pitch-deck-hero-emotional.pptx')
prs.save(out)
print(f"Saved to {out}")
print("""
5 concepts:
  1. "The Invisible Classroom" — 30 students, 12 glowing red, teacher doesn't know
  2. "Every Word Is a Battle" — One giant word with radiating struggle data
  3. "Two Paths" — Diverging futures based on identification
  4. "She's Not Slow. She's Drowning." — Pure typography, scattered letters
  5. "The Book & The Signal" — Physical book + the hidden signal Pacer extracts
""")
