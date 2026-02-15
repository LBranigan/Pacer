#!/usr/bin/env python3
"""Generate an emotional 'crisis' slide with verified stats."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_DEEP = RGBColor(0x06, 0x0C, 0x1A)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
DARK_CARD = RGBColor(0x16, 0x20, 0x32)
ACCENT1 = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
DARKER = RGBColor(0x33, 0x41, 0x55)
RED = RGBColor(0xEF, 0x44, 0x44)
DARK_RED = RGBColor(0xDC, 0x26, 0x26)
DEEP_RED = RGBColor(0xB9, 0x1C, 0x1C)
ORANGE = RGBColor(0xF9, 0x71, 0x16)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x14, 0xB8, 0xA6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txt(slide, text, left, top, width, height, size=18, color=WHITE,
        bold=False, italic=False, align=PP_ALIGN.LEFT, font='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return txBox


def rect(slide, left, top, width, height, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw or 1)
    else: s.line.fill.background()
    return s


def rrect(slide, left, top, width, height, fill=None, line=None, lw=None, r=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = r
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw or 1)
    else: s.line.fill.background()
    return s


def circle(slide, left, top, size, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw or 1)
    else: s.line.fill.background()
    return s


# ============================================================
# THE CRISIS SLIDE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_DEEP)

# ---- TOP SECTION: The headline ----
txt(s, "America has a reading crisis.", Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "And the system designed to catch it is broken.",
    Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.5),
    size=22, color=MUTED, align=PP_ALIGN.CENTER)

# ---- LEFT: THE BIG STATS (NAEP data) ----
# Giant "70%" and "34%" as hero numbers

# --- Card 1: 70% not proficient ---
card1_x = Inches(0.6)
card1_y = Inches(2.2)
card1_w = Inches(3.8)
card1_h = Inches(4.5)

rrect(s, card1_x, card1_y, card1_w, card1_h,
      fill=RGBColor(0x0C, 0x10, 0x1E), line=RGBColor(0x25, 0x15, 0x15), lw=1, r=0.03)

# Red accent bar at top
rect(s, card1_x, card1_y, card1_w, Inches(0.05), fill=RED)

txt(s, "70%", card1_x + Inches(0.3), card1_y + Inches(0.3), Inches(3.2), Inches(1.3),
    size=96, color=RED, bold=True, font='Calibri')

txt(s, "of 8th graders are NOT\nproficient in reading",
    card1_x + Inches(0.3), card1_y + Inches(1.7), Inches(3.2), Inches(0.8),
    size=18, color=SOFT_WHITE, bold=True)

# Visual: bar showing 70/30 split
bar_y = card1_y + Inches(2.7)
# Background bar
rrect(s, card1_x + Inches(0.3), bar_y, Inches(3.2), Inches(0.45),
      fill=RGBColor(0x18, 0x20, 0x30), r=0.15)
# Red portion (70%)
rrect(s, card1_x + Inches(0.3), bar_y, Inches(3.2 * 0.70), Inches(0.45),
      fill=RED, r=0.15)
txt(s, "Not proficient", card1_x + Inches(0.5), bar_y + Inches(0.05), Inches(1.5), Inches(0.35),
    size=10, color=WHITE, bold=True)
txt(s, "30%", card1_x + Inches(2.7), bar_y + Inches(0.05), Inches(0.7), Inches(0.35),
    size=10, color=DIM, bold=True, align=PP_ALIGN.CENTER)

txt(s, "NAEP 2024 \u2014 Lowest scores in\n32 years of testing",
    card1_x + Inches(0.3), card1_y + Inches(3.5), Inches(3.2), Inches(0.7),
    size=11, color=DIM, italic=True)

# --- Card 2: 34% Below Basic ---
card2_x = Inches(4.7)
card2_y = Inches(2.2)
card2_w = Inches(3.8)
card2_h = Inches(4.5)

rrect(s, card2_x, card2_y, card2_w, card2_h,
      fill=RGBColor(0x0C, 0x10, 0x1E), line=RGBColor(0x25, 0x18, 0x08), lw=1, r=0.03)

rect(s, card2_x, card2_y, card2_w, Inches(0.05), fill=ORANGE)

txt(s, "34%", card2_x + Inches(0.3), card2_y + Inches(0.3), Inches(3.2), Inches(1.3),
    size=96, color=ORANGE, bold=True, font='Calibri')

txt(s, "score \"Below Basic\"\nthe highest ever recorded",
    card2_x + Inches(0.3), card2_y + Inches(1.7), Inches(3.2), Inches(0.8),
    size=18, color=SOFT_WHITE, bold=True)

# Bar showing composition
bar2_y = card2_y + Inches(2.7)
rrect(s, card2_x + Inches(0.3), bar2_y, Inches(3.2), Inches(0.45),
      fill=RGBColor(0x18, 0x20, 0x30), r=0.15)
# Below basic (34%)
rrect(s, card2_x + Inches(0.3), bar2_y, Inches(3.2 * 0.34), Inches(0.45),
      fill=DEEP_RED, r=0.15)
# Basic (36%)
rrect(s, card2_x + Inches(0.3 + 3.2 * 0.34), bar2_y, Inches(3.2 * 0.36), Inches(0.45),
      fill=ORANGE)
txt(s, "Below", card2_x + Inches(0.35), bar2_y + Inches(0.05), Inches(0.8), Inches(0.35),
    size=8, color=WHITE, bold=True)
txt(s, "Basic", card2_x + Inches(0.3 + 3.2 * 0.34 + 0.15), bar2_y + Inches(0.05), Inches(0.8), Inches(0.35),
    size=8, color=WHITE, bold=True)
txt(s, "Proficient", card2_x + Inches(2.7), bar2_y + Inches(0.05), Inches(0.8), Inches(0.35),
    size=8, color=DIM, bold=True)

txt(s, "These students can barely decode text,\nlet alone comprehend it",
    card2_x + Inches(0.3), card2_y + Inches(3.5), Inches(3.2), Inches(0.7),
    size=11, color=DIM, italic=True)

# --- Card 3: RTI is broken ---
card3_x = Inches(8.8)
card3_y = Inches(2.2)
card3_w = Inches(3.8)
card3_h = Inches(4.5)

rrect(s, card3_x, card3_y, card3_w, card3_h,
      fill=RGBColor(0x0C, 0x10, 0x1E), line=RGBColor(0x15, 0x18, 0x25), lw=1, r=0.03)

rect(s, card3_x, card3_y, card3_w, Inches(0.05), fill=AMBER)

# The RTI gap
txt(s, "RTI", card3_x + Inches(0.3), card3_y + Inches(0.4), Inches(3.2), Inches(0.5),
    size=20, color=AMBER, bold=True)
txt(s, "The system built to\ncatch these students", card3_x + Inches(0.3), card3_y + Inches(0.8), Inches(3.2), Inches(0.6),
    size=16, color=MUTED)

# Divider
rect(s, card3_x + Inches(0.3), card3_y + Inches(1.55), Inches(3.2), Inches(0.01), fill=DARKER)

# Fidelity stat
txt(s, "20\u201322%", card3_x + Inches(0.3), card3_y + Inches(1.7), Inches(3.2), Inches(0.6),
    size=40, color=AMBER, bold=True, font='Calibri')
txt(s, "implementation fidelity variance",
    card3_x + Inches(0.3), card3_y + Inches(2.25), Inches(3.2), Inches(0.3),
    size=12, color=MUTED, bold=True)

# Divider
rect(s, card3_x + Inches(0.3), card3_y + Inches(2.7), Inches(3.2), Inches(0.01), fill=DARKER)

# The quote - the killer line
txt(s, "\u201c", card3_x + Inches(0.15), card3_y + Inches(2.7), Inches(0.4), Inches(0.5),
    size=36, color=AMBER, font='Georgia')
txt(s, "Teachers report that RTI data collection is so onerous they often just guesstimate to satisfy the system.",
    card3_x + Inches(0.3), card3_y + Inches(3.0), Inches(3.2), Inches(1.2),
    size=12, color=SOFT_WHITE, italic=True, font='Georgia')

# ---- BOTTOM: The connection to Pacer ----
bot_y = Inches(6.85)

# Subtle divider
rect(s, Inches(0.8), bot_y - Inches(0.15), Inches(11.7), Inches(0.01), fill=DARKER)

# Pacer waveform
bx = Inches(4.0)
by = bot_y + Inches(0.0)
for dx, h in [(0, 0.16), (0.09, 0.26), (0.18, 0.38), (0.27, 0.24), (0.36, 0.18)]:
    rrect(s, bx + Inches(dx), by + Inches((0.38-h)/2), Inches(0.055), Inches(h), fill=ACCENT2, r=0.5)
txt(s, "PACER", bx + Inches(0.5), by - Inches(0.02), Inches(1.5), Inches(0.38),
    size=16, color=ACCENT2, bold=True)

txt(s, "Automated struggle detection. Real data. Zero teacher burden.",
    Inches(5.7), bot_y, Inches(5), Inches(0.35),
    size=14, color=MUTED)

# ---- Source citation ----
txt(s, "Source: NAEP 2024 \u2014 The Nation's Report Card  |  nationsreportcard.gov",
    Inches(0.8), bot_y, Inches(3.5), Inches(0.35),
    size=8, color=RGBColor(0x33, 0x41, 0x55), italic=True)


# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pitch-deck-crisis.pptx')
prs.save(out)
print(f"Saved to {out}")
