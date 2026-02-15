"""
Generate Pacer pitch deck v2 as PowerPoint (.pptx)
Dark observatory theme with cyan accents
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ─── Color Palette ───
VOID       = RGBColor(0x03, 0x07, 0x12)
DEEP       = RGBColor(0x06, 0x0C, 0x1A)
PRIMARY    = RGBColor(0x0B, 0x11, 0x20)
SURFACE    = RGBColor(0x11, 0x18, 0x27)
CARD       = RGBColor(0x13, 0x1C, 0x2E)
CARD_LIGHT = RGBColor(0x18, 0x22, 0x38)
CYAN_CORE  = RGBColor(0x06, 0xB6, 0xD4)
CYAN_BRT   = RGBColor(0x22, 0xD3, 0xEE)
CYAN_DEEP  = RGBColor(0x08, 0x91, 0xB2)
SKY        = RGBColor(0x0E, 0xA5, 0xE9)
TEXT_1     = RGBColor(0xF1, 0xF5, 0xF9)
TEXT_2     = RGBColor(0x94, 0xA3, 0xB8)
TEXT_3     = RGBColor(0x64, 0x74, 0x8B)
RED        = RGBColor(0xF4, 0x3F, 0x5E)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
GREEN      = RGBColor(0x34, 0xD3, 0x99)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
TEAL       = RGBColor(0x2D, 0xD4, 0xBF)
BLUE       = RGBColor(0x60, 0xA5, 0xFA)
PINK       = RGBColor(0xF4, 0x72, 0xB6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER      = RGBColor(0xFB, 0xBF, 0x24)
DIM_WHITE  = RGBColor(0xB0, 0xB8, 0xC4)

# ─── Slide dimensions (16:9) ───
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Font names ───
HEADING_FONT = "Outfit"
BODY_FONT    = "DM Sans"
MONO_FONT    = "IBM Plex Mono"

# ─── Helpers ───
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_run(run, text, font_name=BODY_FONT, size=Pt(18), color=TEXT_2, bold=False, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic

def add_heading(slide, text, left, top, width, height=Inches(1.2), size=Pt(40), gradient_text=None):
    """Add heading with optional gradient-colored portion."""
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(4)

    if gradient_text:
        parts = text.split(gradient_text)
        if parts[0]:
            r = p.add_run()
            set_run(r, parts[0], HEADING_FONT, size, TEXT_1, bold=True)
        r = p.add_run()
        set_run(r, gradient_text, HEADING_FONT, size, CYAN_BRT, bold=True)
        if len(parts) > 1 and parts[1]:
            r = p.add_run()
            set_run(r, parts[1], HEADING_FONT, size, TEXT_1, bold=True)
    else:
        r = p.add_run()
        set_run(r, text, HEADING_FONT, size, TEXT_1, bold=True)
    return tb

def add_body_text(slide, text, left, top, width, height, size=Pt(16), color=TEXT_2, bold=False, align=PP_ALIGN.LEFT):
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    set_run(r, text, BODY_FONT, size, color, bold=bold)
    return tb

def add_card(slide, left, top, width, height, fill_color=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Reduce corner radius
    shape.adjustments[0] = 0.04
    return shape

def add_accent_line(slide, left, top, width, color=CYAN_DEEP):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_para(tf, text, size=Pt(15), color=TEXT_2, bold_prefix=None, indent=0):
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.space_after = Pt(2)
    p.level = indent
    if bold_prefix:
        r = p.add_run()
        set_run(r, bold_prefix, BODY_FONT, size, TEXT_1, bold=True)
        r = p.add_run()
        set_run(r, text, BODY_FONT, size, color)
    else:
        r = p.add_run()
        set_run(r, text, BODY_FONT, size, color)
    return p

def add_circle(slide, cx, cy, r, fill_color=None, line_color=None, line_width=Pt(1)):
    left = cx - r
    top = cy - r
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, r*2, r*2)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_arrow_text(slide, left, top, text="→", color=TEXT_3):
    tb = add_textbox(slide, left, top, Inches(0.4), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    set_run(r, text, BODY_FONT, Pt(20), color)
    return tb

def add_pill(slide, left, top, text, dot_color=CYAN_CORE):
    w = Inches(max(1.6, len(text) * 0.11 + 0.5))
    card = add_card(slide, left, top, w, Inches(0.38), CARD)
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left + Inches(0.15), top + Inches(0.12), Inches(0.14), Inches(0.14))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    # Text
    tb = add_textbox(slide, left + Inches(0.35), top + Inches(0.04), w - Inches(0.45), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, text, HEADING_FONT, Pt(11), TEXT_2, bold=True)
    return w


# ═══════════════════════════════════════════
#  BUILD DECK
# ═══════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank layout


# ═══════ SLIDE 1 — HERO ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)

# Left column
L = Inches(0.8)
# Logo
add_body_text(sl, "PACER", L, Inches(1.2), Inches(2.5), Inches(0.5),
              Pt(24), CYAN_BRT, bold=True)
# Waveform bars (simple rectangles as logo)
bar_x = L
bar_y = Inches(1.22)
bar_data = [(0, 0.22), (0.18, 0.35), (0.36, 0.52), (0.54, 0.32), (0.72, 0.2)]
for dx, h in bar_data:
    bh = Inches(h * 0.55)
    by = bar_y + Inches(0.25) - bh/2
    bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        bar_x + Inches(dx * 0.5), by, Inches(0.06), bh)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN_CORE
    bar.line.fill.background()
    bar.adjustments[0] = 0.5

# Headline
tb = add_textbox(sl, L, Inches(1.85), Inches(5.5), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.space_after = Pt(0)
r = p.add_run(); set_run(r, "AI that hears\n", HEADING_FONT, Pt(46), TEXT_1, bold=True)
r = p.add_run(); set_run(r, "how students\nstruggle", HEADING_FONT, Pt(46), CYAN_BRT, bold=True)
r = p.add_run(); set_run(r, " to read", HEADING_FONT, Pt(46), TEXT_1, bold=True)

# Subtext
add_body_text(sl, "Students read aloud from real books. Pacer listens, detects every moment of struggle, and returns rich diagnostic data to teachers\u2014with zero prep.",
    L, Inches(4.05), Inches(5.2), Inches(1.2), Pt(16), TEXT_2)

# Pills
pill_y = Inches(5.4)
px = L
pills = [("Middle School RTI", CYAN_CORE), ("13 Miscue Types", GREEN),
         ("3-Engine ASR", PURPLE), ("Real Books, Not Screens", ORANGE)]
for txt, clr in pills:
    w = add_pill(sl, px, pill_y, txt, clr)
    px += w + Inches(0.12)

# ── Right side: Demo panel ──
R = Inches(7.0)
demo_w = Inches(5.8)
demo_h = Inches(5.6)
demo_y = Inches(0.95)

# Background card for demo
demo_bg = add_card(sl, R, demo_y, demo_w, demo_h, VOID)

# Title bar
title_bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    R, demo_y, demo_w, Inches(0.42))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
title_bar.line.fill.background()

# Traffic light dots
for i, clr in enumerate([RED, AMBER, GREEN]):
    d = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        R + Inches(0.2 + i*0.22), demo_y + Inches(0.14), Inches(0.14), Inches(0.14))
    d.fill.solid()
    d.fill.fore_color.rgb = clr
    d.line.fill.background()

add_body_text(sl, "pacer \u2014 assessment results",
    R + Inches(1.0), demo_y + Inches(0.08), Inches(3), Inches(0.3),
    Pt(9), TEXT_3)

# Demo label
add_body_text(sl, "STUDENT READING \u2014 PASSAGE ANALYSIS",
    R + Inches(0.3), demo_y + Inches(0.55), Inches(4), Inches(0.3),
    Pt(8.5), TEXT_3, bold=True)

# Demo passage with color-coded words
passage_y = demo_y + Inches(0.9)
tb = add_textbox(sl, R + Inches(0.3), passage_y, demo_w - Inches(0.6), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.space_after = Pt(8)
p.line_spacing = Pt(28)

words = [
    ("The ", DIM_WHITE, False, False),
    ("boy ", DIM_WHITE, False, False),
    ("walked ", ORANGE, False, False),
    ("slowly ", DIM_WHITE, False, False),
    ("to ", DIM_WHITE, False, False),
    ("the ", DIM_WHITE, False, False),
    ("gro\u2013groc\u2013grocery ", TEAL, False, False),
    ("store ", DIM_WHITE, False, False),
    ("and ", DIM_WHITE, False, False),
    ("bringed ", ORANGE, False, False),
    ("some ", DIM_WHITE, False, False),
    ("\u22ef 4.1s ", TEXT_3, False, True),
    ("milk ", PURPLE, False, False),
    ("for ", DIM_WHITE, False, False),
    ("his ", DIM_WHITE, False, False),
    ("grandmother ", RED, False, False),
    ("who ", DIM_WHITE, False, False),
    ("lived ", DIM_WHITE, False, False),
    ("around ", ORANGE, False, False),
    ("the ", DIM_WHITE, False, False),
    ("street.", DIM_WHITE, False, False),
]
for txt, clr, bold, is_mono in words:
    r = p.add_run()
    font = MONO_FONT if is_mono else BODY_FONT
    set_run(r, txt, font, Pt(14), clr, bold=bold)

# Annotations above certain words
ann_y = passage_y - Inches(0.08)
add_body_text(sl, "3.2s \u00b7 decoding", R + Inches(2.82), passage_y - Inches(0.12),
    Inches(1.5), Inches(0.2), Pt(7.5), TEAL)
add_body_text(sl, "bought \u2192", R + Inches(4.15), passage_y - Inches(0.12),
    Inches(0.8), Inches(0.2), Pt(7.5), ORANGE)
add_body_text(sl, "self-corrected", R + Inches(1.3), passage_y + Inches(0.52),
    Inches(1.2), Inches(0.2), Pt(7.5), PURPLE)
add_body_text(sl, "across \u2192", R + Inches(3.55), passage_y + Inches(0.52),
    Inches(0.8), Inches(0.2), Pt(7.5), ORANGE)

# Metrics row
met_y = demo_y + Inches(4.05)
# Divider line
div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    R + Inches(0.3), met_y, demo_w - Inches(0.6), Pt(1))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B); div.line.fill.background()

metrics = [
    ("62", "WCPM", CYAN_BRT),
    ("78%", "ACCURACY", CYAN_BRT),
    ("4", "ERRORS", ORANGE),
    ("1", "STRUGGLE", TEAL),
    ("1", "SELF-CORR", PURPLE),
    ("1", "OMISSION", RED),
]
mx = R + Inches(0.3)
for val, lbl, clr in metrics:
    add_body_text(sl, val, mx, met_y + Inches(0.15), Inches(0.8), Inches(0.35),
        Pt(18), clr, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, lbl, mx, met_y + Inches(0.48), Inches(0.8), Inches(0.22),
        Pt(7), TEXT_3, bold=True, align=PP_ALIGN.CENTER)
    mx += Inches(0.9)


# ═══════ SLIDE 2 — PROBLEM ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Reading Fluency Assessment\nIs Broken", Inches(0.8), Inches(0.5),
    Inches(10), Inches(1.4), Pt(40), "Broken")

cards_data = [
    ("Overwhelmed Teachers",
     "120+ students per teacher. Trained to teach content, not diagnose reading struggle. They need help, not more dashboards.",
     RED),
    ("Blunt Metrics",
     "WCPM is like diagnosing heart disease by taking someone's temperature. It misses how a student struggles.",
     ORANGE),
    ("Closing Window",
     "Middle school is the last realistic intervention window. After 8th grade, outcomes calcify. Without identification, there is no intervention.",
     PURPLE),
]
cx = Inches(0.8)
for title, body, accent_clr in cards_data:
    cw = Inches(3.75)
    card = add_card(sl, cx, Inches(2.2), cw, Inches(2.6), CARD)
    add_accent_line(sl, cx, Inches(2.2), cw, accent_clr)
    add_body_text(sl, title, cx + Inches(0.3), Inches(2.55), cw - Inches(0.5), Inches(0.4),
        Pt(17), TEXT_1, bold=True)
    add_body_text(sl, body, cx + Inches(0.3), Inches(3.1), cw - Inches(0.5), Inches(1.4),
        Pt(13), TEXT_2)
    cx += cw + Inches(0.3)

# Funnel
funnel_y = Inches(5.2)
funnel_data = [
    ("Elementary", "Some screening", GREEN, RGBColor(0x0A, 0x2A, 0x1A)),
    ("Middle School", "Gap \u2014 almost nothing", RED, RGBColor(0x2A, 0x0A, 0x12)),
    ("High School", "Too late", TEXT_3, RGBColor(0x18, 0x1C, 0x24)),
]
fx = Inches(2.2)
for i, (title, sub, clr, bg_clr) in enumerate(funnel_data):
    fw = Inches(2.4)
    card = add_card(sl, fx, funnel_y, fw, Inches(0.8), bg_clr)
    add_body_text(sl, title, fx + Inches(0.1), funnel_y + Inches(0.05), fw - Inches(0.2), Inches(0.35),
        Pt(13), clr, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, sub, fx + Inches(0.1), funnel_y + Inches(0.4), fw - Inches(0.2), Inches(0.3),
        Pt(10), TEXT_3, align=PP_ALIGN.CENTER)
    fx += fw
    if i < 2:
        add_arrow_text(sl, fx, funnel_y + Inches(0.15), "\u2192", TEXT_3)
        fx += Inches(0.4)


# ═══════ SLIDE 3 — WHY NOW ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "A LeNet Moment\nfor Speech Recognition", Inches(0.8), Inches(0.4),
    Inches(10), Inches(1.4), Pt(40), "LeNet Moment")

# Left bullets
tb = add_textbox(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0))
tf = tb.text_frame
tf.word_wrap = True
bullets = [
    ("ASR has reached a tipping point\u2014but models are trained on ", "adult, fluent speech"),
    ("Accurate ASR for ", "disfluent populations", " (children, struggling readers, dialect speakers) is the next frontier"),
    ("", "Edge compute", " now makes COPPA/FERPA-compliant classroom deployment practical"),
    ("", "RTI", " is mandated but under-resourced\u2014schools need tools, not more theory"),
]
for parts in bullets:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.space_after = Pt(4)
    p.bullet = True
    for i, part in enumerate(parts):
        r = p.add_run()
        is_bold = (i == 1) if len(parts) > 1 else False
        set_run(r, part, BODY_FONT, Pt(14.5), TEXT_1 if is_bold else TEXT_2, bold=is_bold)

# Right: bar chart
bar_x = Inches(7.2)
# Adult speech bar
add_body_text(sl, "Adult Fluent Speech", bar_x, Inches(2.2), Inches(3), Inches(0.3),
    Pt(12), TEXT_1, bold=True)
add_body_text(sl, "~95%", bar_x + Inches(3.8), Inches(2.2), Inches(1), Inches(0.3),
    Pt(12), CYAN_BRT, bold=True)

track1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    bar_x, Inches(2.55), Inches(5), Inches(0.45))
track1.fill.solid(); track1.fill.fore_color.rgb = CARD; track1.line.fill.background()
track1.adjustments[0] = 0.3

fill1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    bar_x, Inches(2.55), Inches(4.6), Inches(0.45))
fill1.fill.solid(); fill1.fill.fore_color.rgb = CYAN_DEEP; fill1.line.fill.background()
fill1.adjustments[0] = 0.3
add_body_text(sl, "Accurate", bar_x + Inches(0.2), Inches(2.6), Inches(2), Inches(0.35),
    Pt(11), WHITE, bold=True)

# Child speech bar
add_body_text(sl, "Child Disfluent Speech", bar_x, Inches(3.3), Inches(3), Inches(0.3),
    Pt(12), TEXT_1, bold=True)
add_body_text(sl, "~35%", bar_x + Inches(3.8), Inches(3.3), Inches(1), Inches(0.3),
    Pt(12), RED, bold=True)

track2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    bar_x, Inches(3.65), Inches(5), Inches(0.45))
track2.fill.solid(); track2.fill.fore_color.rgb = CARD; track2.line.fill.background()
track2.adjustments[0] = 0.3

fill2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    bar_x, Inches(3.65), Inches(1.7), Inches(0.45))
fill2.fill.solid(); fill2.fill.fore_color.rgb = RED; fill2.line.fill.background()
fill2.adjustments[0] = 0.3
add_body_text(sl, "Unreliable", bar_x + Inches(0.15), Inches(3.7), Inches(1.5), Inches(0.35),
    Pt(11), WHITE, bold=True)

# Gap label
gap_card = add_card(sl, bar_x + Inches(0.4), Inches(4.4), Inches(4.2), Inches(0.55), RGBColor(0x0A, 0x1A, 0x2E))
add_body_text(sl, "\u2195 This gap is the entire opportunity",
    bar_x + Inches(0.5), Inches(4.45), Inches(4), Inches(0.4),
    Pt(14), CYAN_BRT, bold=True, align=PP_ALIGN.CENTER)


# ═══════ SLIDE 4 — SOLUTION ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "A Struggle Detector,\nNot a Score Generator", Inches(0.8), Inches(0.35),
    Inches(10), Inches(1.4), Pt(40), "Struggle Detector")

# Flow steps
flow_data = ["Real Books", "Pacer Listens", "AI Analyzes", "Teacher Sees"]
flow_subs = ["Student reads from\na physical book", "Classroom device\ncaptures audio",
             "Multi-engine struggle\nclassification", "Rich data, zero\nprep required"]
fx = Inches(0.6)
step_w = Inches(2.5)
for i, (title, sub) in enumerate(zip(flow_data, flow_subs)):
    card = add_card(sl, fx, Inches(2.0), step_w, Inches(1.5), CARD)
    add_body_text(sl, title, fx + Inches(0.15), Inches(2.15), step_w - Inches(0.3), Inches(0.4),
        Pt(15), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, sub, fx + Inches(0.15), Inches(2.6), step_w - Inches(0.3), Inches(0.7),
        Pt(11), TEXT_3, align=PP_ALIGN.CENTER)
    fx += step_w
    if i < 3:
        add_arrow_text(sl, fx, Inches(2.45), "\u2192", CYAN_DEEP)
        fx += Inches(0.4)

# Quote box
qbox = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.9))
qbox.fill.solid(); qbox.fill.fore_color.rgb = RGBColor(0x0A, 0x1A, 0x28); qbox.line.fill.background()
qbox.adjustments[0] = 0.08
# Accent left bar
qbar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(3.9), Pt(4), Inches(0.9))
qbar.fill.solid(); qbar.fill.fore_color.rgb = CYAN_DEEP; qbar.line.fill.background()

add_body_text(sl, "\u201CA blood pressure cuff for reading\u2014installed in the classroom, reliable, unobtrusive, always on.\u201D",
    Inches(1.1), Inches(4.0), Inches(11), Inches(0.7), Pt(15), TEXT_1)

# Feature chips
feat_labels = ["Hesitations", "Substitutions", "Repetitions", "Omissions",
               "Self-corrections", "Prosody breakdowns", "Longitudinal trends"]
fx = Inches(0.8)
for lbl in feat_labels:
    chip_w = Inches(max(1.4, len(lbl) * 0.1 + 0.55))
    card = add_card(sl, fx, Inches(5.15), chip_w, Inches(0.4), CARD)
    # Dot
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        fx + Inches(0.12), Inches(5.28), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = CYAN_DEEP; dot.line.fill.background()
    add_body_text(sl, lbl, fx + Inches(0.28), Inches(5.2), chip_w - Inches(0.35), Inches(0.35),
        Pt(11), TEXT_2)
    fx += chip_w + Inches(0.1)


# ═══════ SLIDE 5 — HOW IT WORKS ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "13-Point Miscue Classification", Inches(0.8), Inches(0.3),
    Inches(10), Inches(0.9), Pt(38), "13-Point")

# Pipeline
pipe_steps = ["Audio", "Multi-Engine\nASR (3)", "Forced\nAlignment",
              "Disfluency\nDetection", "Struggle\nClassification", "Dashboard"]
pipe_accent = [False, True, False, True, False, True]
px = Inches(0.5)
step_w = Inches(1.7)
for i, (step, is_acc) in enumerate(zip(pipe_steps, pipe_accent)):
    bg = RGBColor(0x0C, 0x1E, 0x30) if is_acc else CARD
    card = add_card(sl, px, Inches(1.5), step_w, Inches(0.75), bg)
    if is_acc:
        add_accent_line(sl, px, Inches(1.5), step_w, CYAN_BRT)
    add_body_text(sl, step, px + Inches(0.05), Inches(1.55), step_w - Inches(0.1), Inches(0.65),
        Pt(10.5), CYAN_BRT if is_acc else TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    px += step_w
    if i < 5:
        add_arrow_text(sl, px, Inches(1.6), "\u2192", CYAN_DEEP)
        px += Inches(0.35)

# Miscue grid - 3 columns
col_x = [Inches(0.8), Inches(4.6), Inches(8.6)]
col_titles = ["COUNTED AS ERRORS", "DIAGNOSTIC (SHOWN, NOT SCORED)", "FORGIVENESS RULES"]
col_tags = [
    [("Omission", RED), ("Substitution", ORANGE), ("Struggle", TEAL), ("Morphological", ORANGE), ("Long Pause", TEXT_3)],
    [("Insertion", BLUE), ("Hesitation", ORANGE), ("Self-Correction", PURPLE), ("Fragments", PINK), ("Repetitions", PINK), ("Fillers", PURPLE)],
    [("Proper Noun", GREEN)],
]

for ci, (x, title, tags) in enumerate(zip(col_x, col_titles, col_tags)):
    add_body_text(sl, title, x, Inches(2.7), Inches(3.5), Inches(0.3), Pt(8.5), TEXT_3, bold=True)
    ty = Inches(3.1)
    tx = x
    for tag_text, tag_color in tags:
        tag_w = Inches(max(1.1, len(tag_text) * 0.095 + 0.35))
        # Tag background
        tag_bg = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            tx, ty, tag_w, Inches(0.32))
        # Compute bg color (desaturated version)
        r_val = max(0, tag_color[0] // 8)
        g_val = max(0, tag_color[1] // 8)
        b_val = max(0, tag_color[2] // 8)
        tag_bg.fill.solid()
        tag_bg.fill.fore_color.rgb = RGBColor(r_val + 8, g_val + 8, b_val + 12)
        tag_bg.line.color.rgb = tag_color
        tag_bg.line.width = Pt(0.75)
        tag_bg.adjustments[0] = 0.15

        add_body_text(sl, tag_text, tx + Inches(0.08), ty + Inches(0.02),
            tag_w - Inches(0.15), Inches(0.28), Pt(10), tag_color, bold=True, align=PP_ALIGN.CENTER)
        tx += tag_w + Inches(0.08)
        if tx > x + Inches(3.6):
            tx = x
            ty += Inches(0.38)

# Engine badge
add_body_text(sl, "\u26A1 3-Engine Cross-Validation: Reverb \u00b7 Google STT \u00b7 Deepgram Nova-3",
    Inches(0.8), Inches(4.6), Inches(7), Inches(0.4), Pt(12), CYAN_BRT, bold=True)


# ═══════ SLIDE 6 — MARKET ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "$2.4B Market\u2014\nMassive Gap in the Middle", Inches(0.8), Inches(0.4),
    Inches(10), Inches(1.4), Pt(40), "$2.4B")

# Concentric circles (left)
ccx = Inches(3.2)
ccy = Inches(4.2)
# Outer
add_circle(sl, ccx, ccy, Inches(1.9), RGBColor(0x0D, 0x13, 0x22), TEXT_3, Pt(1))
add_body_text(sl, "Expansion: SLP, Private, Homeschool",
    ccx - Inches(1.5), ccy - Inches(1.7), Inches(3), Inches(0.3),
    Pt(8.5), TEXT_3, bold=True, align=PP_ALIGN.CENTER)
# Mid
add_circle(sl, ccx, ccy, Inches(1.3), RGBColor(0x09, 0x15, 0x28), SKY, Pt(1))
add_body_text(sl, "Adjacent: K-5, High School",
    ccx - Inches(1.1), ccy - Inches(1.1), Inches(2.2), Inches(0.3),
    Pt(8.5), SKY, bold=True, align=PP_ALIGN.CENTER)
# Core
add_circle(sl, ccx, ccy, Inches(0.72), RGBColor(0x08, 0x1C, 0x2C), CYAN_CORE, Pt(2))
add_body_text(sl, "Core:\nMiddle School\nRTI Screening",
    ccx - Inches(0.6), ccy - Inches(0.4), Inches(1.2), Inches(0.8),
    Pt(9), TEXT_1, bold=True, align=PP_ALIGN.CENTER)

# Right: market points
tb = add_textbox(sl, Inches(6.5), Inches(2.2), Inches(6), Inches(5))
tf = tb.text_frame
tf.word_wrap = True
market_bullets = [
    ("Middle school", " is drastically underserved\u2014less money, more complex issues, fewer competitors"),
    ("B2B entry:", " school districts via RTI/MTSS mandates (existing budget line items)"),
    ("B2C expansion:", " SLPs, private schools, homeschool families"),
    ("Data moat:", " every session = proprietary disfluent children's speech data\u2014the scarcest, most valuable dataset in edtech ASR"),
]
for bold_part, rest in market_bullets:
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.space_after = Pt(4)
    p.bullet = True
    r = p.add_run(); set_run(r, bold_part, BODY_FONT, Pt(15), TEXT_1, bold=True)
    r = p.add_run(); set_run(r, rest, BODY_FONT, Pt(15), TEXT_2)


# ═══════ SLIDE 7 — COMPETITION ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Positioned Where\nNo One Else Is Playing", Inches(0.8), Inches(0.3),
    Inches(10), Inches(1.4), Pt(40), "No One Else Is Playing")

# 2x2 Matrix
mx = Inches(2.5)
my = Inches(2.0)
mw = Inches(4.0)
mh = Inches(2.0)
gap = Inches(0.06)

cells = [
    # (col, row, names, is_highlight)
    (0, 0, [("NWEA MAP", "Legacy \u00b7 Standardized scores"), ("iStation", "Legacy \u00b7 WCPM only")], False),
    (1, 0, [], False),  # empty
    (0, 1, [("Amira Learning", "~$40M raised \u00b7 AI tutoring"), ("Ello", "Reed Hastings, YC \u00b7 K-2")], False),
    (1, 1, [("PACER", "Alone in this quadrant")], True),  # highlight
]

for col, row, names, hl in cells:
    cx = mx + col * (mw + gap)
    cy = my + row * (mh + gap)
    bg = RGBColor(0x0C, 0x1E, 0x30) if hl else CARD
    cell_shape = add_card(sl, cx, cy, mw, mh, bg)

    if hl:
        # PACER badge
        add_body_text(sl, "PACER", cx + Inches(0.2), cy + Inches(0.5), mw - Inches(0.4), Inches(0.6),
            Pt(28), CYAN_BRT, bold=True, align=PP_ALIGN.CENTER)
        add_body_text(sl, "Alone in this quadrant", cx + Inches(0.2), cy + Inches(1.15), mw - Inches(0.4), Inches(0.4),
            Pt(11), CYAN_DEEP, bold=True, align=PP_ALIGN.CENTER)
    elif not names:
        add_body_text(sl, "\u2014", cx + Inches(0.2), cy + Inches(0.7), mw - Inches(0.4), Inches(0.4),
            Pt(14), TEXT_3, align=PP_ALIGN.CENTER)
    else:
        ny = cy + Inches(0.25)
        for name, sub in names:
            add_body_text(sl, name, cx + Inches(0.2), ny, mw - Inches(0.4), Inches(0.35),
                Pt(13), TEXT_2, bold=True, align=PP_ALIGN.CENTER)
            add_body_text(sl, sub, cx + Inches(0.2), ny + Inches(0.3), mw - Inches(0.4), Inches(0.3),
                Pt(9), TEXT_3, align=PP_ALIGN.CENTER)
            ny += Inches(0.7)

# Axis labels
add_body_text(sl, "Score-level  \u2192  Struggle-level granularity",
    Inches(3.5), Inches(6.2), Inches(5), Inches(0.3), Pt(9), TEXT_3, bold=True, align=PP_ALIGN.CENTER)
add_body_text(sl, "SCREEN-BASED", mx - Inches(1.8), Inches(2.8), Inches(1.5), Inches(0.3),
    Pt(8), TEXT_3, bold=True, align=PP_ALIGN.CENTER)
add_body_text(sl, "PHYSICAL BOOKS", mx - Inches(1.8), Inches(4.8), Inches(1.5), Inches(0.3),
    Pt(8), TEXT_3, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
add_body_text(sl, "Amira requires talking to a computer. Ello targets K-2. Legacy tools give scores, not insights.",
    Inches(1.5), Inches(6.6), Inches(10), Inches(0.4), Pt(12), TEXT_3, align=PP_ALIGN.CENTER)


# ═══════ SLIDE 8 — TRACTION ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Already in Classrooms", Inches(0.8), Inches(0.35),
    Inches(10), Inches(0.9), Pt(40), "Already")

# Timeline
tl_data = [
    ("\u2713", "Built MVP", "Full AI pipeline live", False),
    ("\u25C9", "Morningside Pilot", "Active\u2014real students", True),
    ("3-5", "School Expansion", "Precision teaching network", False),
    ("\u2192", "District Partners", "Paid pilots", False),
]
tx_start = Inches(1.5)
tl_step_w = Inches(2.5)
# Timeline line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    tx_start + Inches(0.5), Inches(1.95), Inches(8.5), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = SURFACE; line.line.fill.background()

for i, (icon, title, sub, active) in enumerate(tl_data):
    cx = tx_start + i * tl_step_w
    # Dot
    d_size = Inches(0.42)
    d = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        cx + Inches(0.8), Inches(1.72), d_size, d_size)
    d.fill.solid()
    d.fill.fore_color.rgb = RGBColor(0x08, 0x22, 0x35) if active else CARD
    d.line.color.rgb = CYAN_CORE if active else SURFACE
    d.line.width = Pt(2)

    add_body_text(sl, icon, cx + Inches(0.8), Inches(1.75), d_size, d_size,
        Pt(12), CYAN_BRT if active else TEXT_2, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, title, cx + Inches(0.1), Inches(2.3), Inches(1.8), Inches(0.35),
        Pt(12), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, sub, cx + Inches(0.1), Inches(2.6), Inches(1.8), Inches(0.3),
        Pt(10), TEXT_3, align=PP_ALIGN.CENTER)

# Traction cards
tract_data = [
    ("Morningside Academy", "\u2014live pilot, real students generating real data, zero bureaucratic friction"),
    ("Andrew Kieta", "\u2014precision teaching network opens doors to schools already committed to data-driven reading"),
    ("Product functional today", "\u2014Word Speed Map, disfluency detection, multi-miscue engine all live"),
    ("COPPA/FERPA compliant", " architecture from day one\u2014edge compute, no cloud-stored student audio"),
]
tcx = Inches(0.8)
tcy = Inches(3.3)
for i, (bold_part, rest) in enumerate(tract_data):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.0)
    y = Inches(3.3) + row * Inches(1.3)
    card = add_card(sl, x, y, Inches(5.7), Inches(1.1), CARD)

    tb = add_textbox(sl, x + Inches(0.25), y + Inches(0.2), Inches(5.2), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, bold_part, BODY_FONT, Pt(13), TEXT_1, bold=True)
    r = p.add_run(); set_run(r, rest, BODY_FONT, Pt(13), TEXT_2)


# ═══════ SLIDE 9 — TEAM ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Built for This Problem", Inches(0.8), Inches(0.35),
    Inches(10), Inches(0.9), Pt(40), "Built")

# Team cards
team_data = [
    ("Emma", [
        ("PhD from ", "Vanderbilt", "\u2014#1 Special Education program"),
        ("Mentored by ", "Doug Fuchs", " (creator of PALs system)"),
        ("Dual expertise: ", "Linguistics + Special Education", ""),
        ("Decade of classroom teaching experience", "", ""),
        ("Currently runs the ", "Wing Institute", ""),
    ]),
    ("Founder", [
        ("", "10+ years", " as entrepreneur"),
        ("Built Pacer's ", "entire AI pipeline", " solo"),
        ("All-in: ", "12+ hour days", ", 7 days/week"),
        ("Deep technical background\u2014AI/ML, edge compute, audio", "", ""),
        ("Former Amazon seller\u2014understands ", "product-market fit", ""),
    ]),
]

for ci, (role, creds) in enumerate(team_data):
    cx = Inches(0.8) + ci * Inches(6.0)
    cw = Inches(5.7)
    card = add_card(sl, cx, Inches(1.5), cw, Inches(3.5), CARD)
    add_accent_line(sl, cx, Inches(1.5), cw, CYAN_BRT)

    add_body_text(sl, role, cx + Inches(0.3), Inches(1.75), cw - Inches(0.5), Inches(0.45),
        Pt(20), TEXT_1, bold=True)

    tb = add_textbox(sl, cx + Inches(0.3), Inches(2.3), cw - Inches(0.5), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    for parts in creds:
        p = tf.add_paragraph()
        p.space_before = Pt(3); p.space_after = Pt(1)
        p.bullet = True
        for j, part in enumerate(parts):
            if not part: continue
            r = p.add_run()
            is_bold = (j == 1) and part != ""
            set_run(r, part, BODY_FONT, Pt(12.5), TEXT_1 if is_bold else TEXT_2, bold=is_bold)

# Network bar
net_data = [
    ("Morningside Academy", "Free classroom access \u00b7 Pilot partner \u00b7 Network gateway"),
    ("Nvidia Connection", "Connor\u2014Distinguished Engineer \u00b7 GPU/inference expertise"),
]
for ni, (title, sub) in enumerate(net_data):
    nx = Inches(0.8) + ni * Inches(6.0)
    card = add_card(sl, nx, Inches(5.3), Inches(5.7), Inches(0.85), CARD)
    add_body_text(sl, title, nx + Inches(0.25), Inches(5.38), Inches(5), Inches(0.35),
        Pt(12.5), TEXT_1, bold=True)
    add_body_text(sl, sub, nx + Inches(0.25), Inches(5.7), Inches(5), Inches(0.3),
        Pt(10), TEXT_3)


# ═══════ SLIDE 10 — BUSINESS MODEL ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Land with Screening,\nExpand with Data", Inches(0.8), Inches(0.35),
    Inches(10), Inches(1.4), Pt(40), "Expand with Data")

bm_data = [
    ("01 / LAND", "RTI Screening", "Per-school SaaS subscription.\nAligns with existing budget\nline items and mandates.", True),
    ("02 / EXPAND", "District Analytics", "Progress monitoring, longitudinal\ndashboards, tier 2/3\nintervention tracking.", False),
    ("03 / DEFEND", "Data Moat", "Proprietary disfluent speech data.\nLicense fine-tuned model weights\nto other platforms.", False),
]
for i, (num, title, desc, primary) in enumerate(bm_data):
    bx = Inches(0.8) + i * Inches(4.0)
    bw = Inches(3.7)
    bg = RGBColor(0x0A, 0x1C, 0x2E) if primary else CARD
    card = add_card(sl, bx, Inches(2.1), bw, Inches(2.5), bg)
    if primary:
        add_accent_line(sl, bx, Inches(2.1), bw, CYAN_BRT)

    add_body_text(sl, num, bx + Inches(0.25), Inches(2.3), bw - Inches(0.4), Inches(0.3),
        Pt(10), CYAN_DEEP, bold=True)
    add_body_text(sl, title, bx + Inches(0.25), Inches(2.65), bw - Inches(0.4), Inches(0.4),
        Pt(17), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, desc, bx + Inches(0.25), Inches(3.15), bw - Inches(0.4), Inches(1.2),
        Pt(12), TEXT_2, align=PP_ALIGN.CENTER)

# Callout
callout_bg = add_card(sl, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.3), RGBColor(0x18, 0x0C, 0x10))
tb = add_textbox(sl, Inches(1.2), Inches(5.15), Inches(10.8), Inches(1.0))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); set_run(r, "Soapbox Labs", BODY_FONT, Pt(13), TEXT_1, bold=True)
r = p.add_run(); set_run(r, " was acquired by Curriculum Associates for an estimated ", BODY_FONT, Pt(13), TEXT_2)
r = p.add_run(); set_run(r, "$100M+", BODY_FONT, Pt(13), TEXT_1, bold=True)
r = p.add_run(); set_run(r, "\u2014for the children's speech data. First thing they did: closed the API to competitors. Our data will be richer because we capture ", BODY_FONT, Pt(13), TEXT_2)
r = p.add_run(); set_run(r, "disfluency", BODY_FONT, Pt(13), TEXT_1, bold=True, italic=True)
r = p.add_run(); set_run(r, ", not just correct speech.", BODY_FONT, Pt(13), TEXT_2)


# ═══════ SLIDE 11 — THE ASK ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)
add_heading(sl, "Accelerate Pacer into\n50 Classrooms", Inches(0.8), Inches(0.3),
    Inches(10), Inches(1.4), Pt(40), "50 Classrooms")

add_body_text(sl, "Seeking Launch.co Partnership",
    Inches(0.8), Inches(1.85), Inches(11.5), Inches(0.45),
    Pt(20), TEXT_2, align=PP_ALIGN.CENTER)

ask_data = [
    ("Hire", "First ML engineer to accelerate\nmodel development"),
    ("Deploy", "Edge hardware in classrooms\nfor on-device inference"),
    ("Validate", "3 paid district pilots to\nprove sales motion"),
]
for i, (title, desc) in enumerate(ask_data):
    ax = Inches(0.8) + i * Inches(4.0)
    aw = Inches(3.7)
    card = add_card(sl, ax, Inches(2.6), aw, Inches(1.8), CARD)
    add_body_text(sl, title, ax + Inches(0.2), Inches(2.8), aw - Inches(0.4), Inches(0.4),
        Pt(17), TEXT_1, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, desc, ax + Inches(0.2), Inches(3.3), aw - Inches(0.4), Inches(0.9),
        Pt(12), TEXT_2, align=PP_ALIGN.CENTER)

# Milestones bar
mst_bg = add_card(sl, Inches(1.5), Inches(4.9), Inches(10.3), Inches(1.2), CARD)
mst_data = [("50", "Classrooms"), ("3", "District LOIs"), ("1", "Validation Study"), ("12 mo", "Timeline")]
for i, (val, lbl) in enumerate(mst_data):
    mx = Inches(1.8) + i * Inches(2.5)
    add_body_text(sl, val, mx, Inches(5.0), Inches(2), Inches(0.5),
        Pt(26) if val != "12 mo" else Pt(16), CYAN_BRT, bold=True, align=PP_ALIGN.CENTER)
    add_body_text(sl, lbl, mx, Inches(5.5), Inches(2), Inches(0.3),
        Pt(10), TEXT_3, bold=True, align=PP_ALIGN.CENTER)
    # Separator
    if i < 3:
        sep = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            mx + Inches(2.1), Inches(5.35), Inches(0.3), Pt(1))
        sep.fill.solid(); sep.fill.fore_color.rgb = SURFACE; sep.line.fill.background()


# ═══════ SLIDE 12 — CLOSING ═══════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, PRIMARY)

# Logo wordmark
add_body_text(sl, "PACER", Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.6),
    Pt(18), TEXT_3, bold=True, align=PP_ALIGN.CENTER)

# Waveform bars centered
bar_center_x = Inches(6.2)
bar_y_c = Inches(1.65)
for dx, h in bar_data:
    bh = Inches(h * 0.45)
    by = bar_y_c + Inches(0.2) - bh/2
    bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        bar_center_x + Inches(dx * 0.45), by, Inches(0.055), bh)
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN_CORE
    bar.line.fill.background()
    bar.adjustments[0] = 0.5

# Main headline
tb = add_textbox(sl, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); set_run(r, "Every Struggling Reader ", HEADING_FONT, Pt(48), TEXT_1, bold=True)
r = p.add_run(); set_run(r, "Found", HEADING_FONT, Pt(48), CYAN_BRT, bold=True)
r = p.add_run(); set_run(r, ".", HEADING_FONT, Pt(48), TEXT_1, bold=True)

add_body_text(sl, "None Left Behind.", Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7),
    Pt(28), TEXT_2, align=PP_ALIGN.CENTER)

add_body_text(sl, "Pacer \u2014 AI-Powered Reading Struggle Detection",
    Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.4),
    Pt(13), TEXT_3, align=PP_ALIGN.CENTER)


# ═══════ SAVE ═══════
output_path = "/mnt/c/Users/brani/desktop/googstt/pitch-deck-v2.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
