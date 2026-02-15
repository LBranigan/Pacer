#!/usr/bin/env python3
"""Generate 3 hero slide concepts for Pacer pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import os

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_DEEP = RGBColor(0x06, 0x0C, 0x1A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
DARK_CARD = RGBColor(0x16, 0x20, 0x32)
ACCENT1 = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
BRIGHT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
DARKER_DIM = RGBColor(0x47, 0x55, 0x69)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x71, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WARM_BG = RGBColor(0xF5, 0xF0, 0xE8)
CREAM = RGBColor(0xFA, 0xF7, 0xF2)
BOOK_BROWN = RGBColor(0x8B, 0x6F, 0x52)
BOOK_SPINE = RGBColor(0x6B, 0x52, 0x3C)
BOOK_PAGE = RGBColor(0xFD, 0xFA, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE,
             bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri',
             anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    # Set vertical anchor
    txBox.text_frame._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 't'))
    run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height, font_name='Calibri', alignment=PP_ALIGN.LEFT):
    """lines = list of (text, font_size, color, bold, italic)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, italic) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
    return txBox


def shape(slide, stype, left, top, width, height, fill_color=None, line_color=None, line_w=None):
    s = slide.shapes.add_shape(stype, left, top, width, height)
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        if line_w:
            s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s


def rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_w=None, radius=0.05):
    s = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color, line_color, line_w)
    s.adjustments[0] = radius
    return s


def add_gradient_bg(slide, color1, color2):
    """Add a gradient background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


# ============================================================
# CONCEPT A: "BEFORE / AFTER" — What teachers see vs what Pacer reveals
# Bold split-screen. Left is bleak, right is rich.
# ============================================================
sA = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sA, DARK_BG)

# --- Top branding ---
# Waveform bars
bx = Inches(0.7)
by = Inches(0.45)
for dx, h, op_color in [(0, 0.18, ACCENT1), (0.1, 0.3, ACCENT1), (0.2, 0.45, ACCENT2), (0.3, 0.28, ACCENT1), (0.4, 0.2, ACCENT1)]:
    bar = rounded_rect(sA, bx + Inches(dx), by + Inches((0.45-h)/2), Inches(0.065), Inches(h), fill_color=op_color, radius=0.5)
add_text(sA, "PACER", Inches(1.35), Inches(0.35), Inches(2), Inches(0.5),
         font_size=22, color=ACCENT2, bold=True)
add_text(sA, "AI-Powered Reading Struggle Detection", Inches(1.35), Inches(0.68), Inches(4), Inches(0.3),
         font_size=11, color=DIM)
add_text(sA, "CONCEPT A", Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.3),
         font_size=10, color=DIM, bold=True, alignment=PP_ALIGN.RIGHT)

# --- Main headline ---
add_text(sA, "See what you've been missing.", Inches(0.7), Inches(1.2), Inches(12), Inches(0.8),
         font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# --- LEFT PANEL: "What teachers see today" ---
left_x = Inches(0.7)
panel_y = Inches(2.4)
panel_w = Inches(5.7)
panel_h = Inches(4.5)

rounded_rect(sA, left_x, panel_y, panel_w, panel_h,
             fill_color=RGBColor(0x0A, 0x10, 0x1F), line_color=RGBColor(0x1E, 0x29, 0x3B), line_w=1, radius=0.03)

add_text(sA, "WHAT TEACHERS SEE TODAY", left_x + Inches(0.4), panel_y + Inches(0.3),
         Inches(4), Inches(0.3), font_size=10, color=DIM, bold=True)

# Big lonely WCPM number
add_text(sA, "87", Inches(2.0), panel_y + Inches(1.2), Inches(2.8), Inches(1.6),
         font_size=96, color=RGBColor(0x33, 0x41, 0x55), bold=True, alignment=PP_ALIGN.CENTER)
add_text(sA, "WCPM", Inches(2.0), panel_y + Inches(2.7), Inches(2.8), Inches(0.4),
         font_size=18, color=RGBColor(0x33, 0x41, 0x55), alignment=PP_ALIGN.CENTER)

add_text(sA, "A single number.\nNo context. No diagnosis. No action.", Inches(1.5), panel_y + Inches(3.4), Inches(4), Inches(0.8),
         font_size=14, color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

# --- RIGHT PANEL: "What Pacer reveals" ---
right_x = Inches(6.9)

rounded_rect(sA, right_x, panel_y, panel_w, panel_h,
             fill_color=RGBColor(0x0A, 0x18, 0x28), line_color=ACCENT1, line_w=1.5, radius=0.03)

add_text(sA, "WHAT PACER REVEALS", right_x + Inches(0.4), panel_y + Inches(0.3),
         Inches(4), Inches(0.3), font_size=10, color=ACCENT2, bold=True)

# Simulated color-coded passage - using a simple text block approach
# Row of "words" as colored rounded rects
passage_y = panel_y + Inches(0.8)
row_words = [
    # Row 1
    [("The", MUTED, None), ("boy", MUTED, None), ("walked", ORANGE, "hesitation"),
     ("slowly", MUTED, None), ("to", MUTED, None), ("the", MUTED, None)],
    # Row 2
    [("grocery", TEAL, "struggle \u2022 3.2s"), ("store", MUTED, None), ("and", MUTED, None),
     ("bringed", ORANGE, "substitution"), ("some", MUTED, None)],
    # Row 3
    [("\u23f8 4.1s", GRAY, "long pause"), ("milk", PURPLE, "self-corrected"),
     ("for", MUTED, None), ("his", MUTED, None), ("grandmother", RED, "omission")],
]

for row_i, row in enumerate(row_words):
    wy = passage_y + Inches(row_i * 0.95)
    wx = right_x + Inches(0.4)
    for word, color, annotation in row:
        ww = max(len(word) * 0.12 + 0.15, 0.55)

        if annotation:
            # Annotation above
            add_text(sA, annotation, wx, wy - Inches(0.22), Inches(ww + 0.2), Inches(0.2),
                     font_size=7, color=color, font_name='Consolas')

        # Word chip
        chip_color = RGBColor(0x12, 0x22, 0x35) if color == MUTED else None
        if color != MUTED:
            # Colored background chip
            bg_colors = {
                ORANGE: RGBColor(0x2A, 0x1A, 0x0A),
                TEAL: RGBColor(0x0A, 0x22, 0x1E),
                PURPLE: RGBColor(0x1A, 0x10, 0x28),
                RED: RGBColor(0x2A, 0x0A, 0x0A),
                GRAY: RGBColor(0x18, 0x1C, 0x22),
            }
            chip_bg = bg_colors.get(color, RGBColor(0x12, 0x22, 0x35))
            rounded_rect(sA, wx, wy, Inches(ww), Inches(0.38),
                         fill_color=chip_bg, line_color=color, line_w=1, radius=0.15)
        else:
            rounded_rect(sA, wx, wy, Inches(ww), Inches(0.38),
                         fill_color=RGBColor(0x12, 0x22, 0x35), line_color=RGBColor(0x25, 0x30, 0x45), line_w=0.5, radius=0.15)

        add_text(sA, word, wx + Inches(0.08), wy + Inches(0.02), Inches(ww - 0.1), Inches(0.34),
                 font_size=13, color=color if color != MUTED else RGBColor(0x8A, 0x94, 0xA5),
                 bold=(color != MUTED), font_name='Calibri')

        wx += Inches(ww + 0.1)

# Metrics row at bottom of right panel
met_y = panel_y + Inches(3.6)
shape(sA, MSO_SHAPE.RECTANGLE, right_x + Inches(0.4), met_y, Inches(4.9), Inches(0.015),
      fill_color=RGBColor(0x1E, 0x29, 0x3B))

metrics = [("62", "WCPM", ACCENT2), ("78%", "Accuracy", ACCENT2), ("4", "Errors", ORANGE),
           ("1", "Struggle", TEAL), ("1", "Self-Corr", PURPLE), ("1", "Omission", RED)]
for i, (val, label, color) in enumerate(metrics):
    mx = right_x + Inches(0.4 + i * 0.87)
    add_text(sA, val, mx, met_y + Inches(0.1), Inches(0.8), Inches(0.35),
             font_size=20, color=color, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
    add_text(sA, label, mx, met_y + Inches(0.45), Inches(0.8), Inches(0.2),
             font_size=8, color=DIM, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# CONCEPT B: "BOOK → INSIGHT" — Physical book focus with flow
# Warm left side (book world) to cool right side (data world)
# ============================================================
sB = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sB, DARK_BG)

add_text(sB, "CONCEPT B", Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.3),
         font_size=10, color=DIM, bold=True, alignment=PP_ALIGN.RIGHT)

# --- LEFT THIRD: The Physical Book World ---
# Warm-toned panel
book_panel_w = Inches(4.5)
rounded_rect(sB, Inches(0.4), Inches(0.4), book_panel_w, Inches(6.7),
             fill_color=RGBColor(0x1C, 0x19, 0x14), line_color=RGBColor(0x3A, 0x32, 0x28), line_w=1, radius=0.02)

# Big physical book illustration
# Book cover (front)
bk_x = Inches(1.1)
bk_y = Inches(1.4)
bk_w = Inches(2.8)
bk_h = Inches(3.5)
# Shadow
rounded_rect(sB, bk_x + Inches(0.08), bk_y + Inches(0.08), bk_w, bk_h,
             fill_color=RGBColor(0x0A, 0x08, 0x06), radius=0.02)
# Pages (visible edge)
rounded_rect(sB, bk_x + Inches(0.12), bk_y + Inches(0.08), bk_w - Inches(0.12), bk_h - Inches(0.08),
             fill_color=BOOK_PAGE, radius=0.01)
# Book cover
rounded_rect(sB, bk_x, bk_y, bk_w, bk_h,
             fill_color=BOOK_BROWN, line_color=BOOK_SPINE, line_w=1, radius=0.02)
# Spine
shape(sB, MSO_SHAPE.RECTANGLE, bk_x, bk_y, Inches(0.18), bk_h,
      fill_color=BOOK_SPINE)
# Title area on cover
rounded_rect(sB, bk_x + Inches(0.5), bk_y + Inches(0.8), Inches(1.8), Inches(0.6),
             fill_color=RGBColor(0x7A, 0x5F, 0x44), line_color=RGBColor(0x9A, 0x80, 0x65), line_w=0.5, radius=0.03)
add_text(sB, "Island of the\nBlue Dolphins", bk_x + Inches(0.55), bk_y + Inches(0.82), Inches(1.7), Inches(0.55),
         font_size=11, color=CREAM, bold=True, alignment=PP_ALIGN.CENTER, font_name='Georgia')
# Author
add_text(sB, "Scott O'Dell", bk_x + Inches(0.5), bk_y + Inches(1.55), Inches(1.8), Inches(0.3),
         font_size=9, color=RGBColor(0xC4, 0xAE, 0x96), alignment=PP_ALIGN.CENTER, font_name='Georgia', italic=True)

# "Grade-level text" label
add_text(sB, "Real books. Grade-appropriate text.\nNot a screen.", Inches(0.7), Inches(5.3), Inches(3.8), Inches(0.7),
         font_size=14, color=RGBColor(0xB0, 0x9A, 0x80), alignment=PP_ALIGN.CENTER, italic=True)

# Label below
add_text(sB, "STUDENT READS ALOUD", Inches(0.7), Inches(6.2), Inches(3.8), Inches(0.3),
         font_size=10, color=RGBColor(0x6B, 0x5E, 0x50), bold=True, alignment=PP_ALIGN.CENTER)

# --- MIDDLE: The Transition (waveform / audio) ---
mid_x = Inches(5.2)

# Large arrow / waveform zone
add_text(sB, "\u25B6", Inches(4.95), Inches(3.2), Inches(0.5), Inches(0.5),
         font_size=28, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Waveform visualization
wave_y = Inches(2.0)
add_text(sB, "PACER LISTENS", mid_x + Inches(0.1), Inches(1.4), Inches(2.8), Inches(0.3),
         font_size=10, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

# Waveform bars (audio visualization)
wave_cx = mid_x + Inches(1.5)
wave_cy = Inches(3.2)
bar_heights = [0.3, 0.5, 0.8, 1.2, 0.9, 1.4, 0.7, 1.1, 0.6, 0.9, 1.3, 0.5, 0.8, 0.4, 0.6, 1.0, 0.7, 0.3, 0.5, 0.2]
for i, bh in enumerate(bar_heights):
    bar_color = ACCENT2 if i in [3, 5, 10] else ACCENT1  # highlight peaks
    bx_pos = wave_cx - Inches(len(bar_heights) * 0.065) + Inches(i * 0.13)
    rounded_rect(sB, bx_pos, wave_cy - Inches(bh * 0.4),
                 Inches(0.08), Inches(bh * 0.8),
                 fill_color=bar_color, radius=0.5)

# AI analysis labels in middle
analysis_items = [
    ("Hesitation Detection", ORANGE),
    ("Substitution Analysis", ORANGE),
    ("Struggle Classification", TEAL),
    ("Prosody Tracking", PURPLE),
    ("Omission Detection", RED),
]
for i, (label, color) in enumerate(analysis_items):
    ay = Inches(4.2 + i * 0.38)
    rounded_rect(sB, mid_x + Inches(0.3), ay, Inches(2.3), Inches(0.3),
                 fill_color=RGBColor(0x0A, 0x14, 0x25), line_color=color, line_w=0.5, radius=0.4)
    # small dot
    shape(sB, MSO_SHAPE.OVAL, mid_x + Inches(0.42), ay + Inches(0.09),
          Inches(0.1), Inches(0.1), fill_color=color)
    add_text(sB, label, mid_x + Inches(0.6), ay + Inches(0.01), Inches(1.9), Inches(0.28),
             font_size=9, color=color, bold=True)

add_text(sB, "3-ENGINE ASR", mid_x + Inches(0.3), Inches(6.2), Inches(2.3), Inches(0.3),
         font_size=10, color=DIM, bold=True, alignment=PP_ALIGN.CENTER)

# --- RIGHT: Arrow to insights ---
add_text(sB, "\u25B6", Inches(8.15), Inches(3.2), Inches(0.5), Inches(0.5),
         font_size=28, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Right panel: Teacher insight
right_x = Inches(8.6)
right_w = Inches(4.3)
rounded_rect(sB, right_x, Inches(0.4), right_w, Inches(6.7),
             fill_color=RGBColor(0x0A, 0x14, 0x25), line_color=ACCENT1, line_w=1, radius=0.02)

add_text(sB, "TEACHER RECEIVES", right_x + Inches(0.3), Inches(0.6), Inches(3.5), Inches(0.3),
         font_size=10, color=ACCENT2, bold=True)

# Metric cards
metric_cards = [
    ("62 WCPM", "Reading Rate", ACCENT2),
    ("78%", "Accuracy", ACCENT2),
    ("4 errors", "Identified", ORANGE),
    ("1 struggle", "Flagged", TEAL),
]
for i, (val, label, color) in enumerate(metric_cards):
    row = i // 2
    col = i % 2
    cx = right_x + Inches(0.3 + col * 2.0)
    cy = Inches(1.1 + row * 1.2)
    rounded_rect(sB, cx, cy, Inches(1.8), Inches(1.0),
                 fill_color=RGBColor(0x10, 0x1A, 0x2E), line_color=RGBColor(0x25, 0x30, 0x45), line_w=0.75, radius=0.06)
    add_text(sB, val, cx + Inches(0.15), cy + Inches(0.15), Inches(1.5), Inches(0.4),
             font_size=22, color=color, bold=True, font_name='Consolas')
    add_text(sB, label, cx + Inches(0.15), cy + Inches(0.6), Inches(1.5), Inches(0.25),
             font_size=10, color=DIM)

# Longitudinal insight
add_text(sB, "LONGITUDINAL TREND", right_x + Inches(0.3), Inches(3.6), Inches(3.5), Inches(0.25),
         font_size=9, color=DIM, bold=True)

# Simple trend line (rising)
trend_y_base = Inches(4.9)
trend_points = [0.8, 0.7, 0.75, 0.55, 0.5, 0.35, 0.3, 0.2]
for i, ty in enumerate(trend_points):
    dot_x = right_x + Inches(0.5 + i * 0.45)
    dot_y = Inches(3.9) + Inches(ty)
    dot_color = RED if i < 3 else (ORANGE if i < 5 else GREEN)
    shape(sB, MSO_SHAPE.OVAL, dot_x, dot_y, Inches(0.12), Inches(0.12), fill_color=dot_color)
    if i < len(trend_points) - 1:
        # Connector line (horizontal)
        next_y = Inches(3.9) + Inches(trend_points[i+1])
        shape(sB, MSO_SHAPE.RECTANGLE, dot_x + Inches(0.1), min(dot_y, next_y) + Inches(0.05),
              Inches(0.37), Inches(0.02), fill_color=RGBColor(0x25, 0x30, 0x45))

add_text(sB, "Sessions \u2192", right_x + Inches(0.5), Inches(5.1), Inches(3), Inches(0.2),
         font_size=8, color=DIM)
add_text(sB, "Improving \u2191", right_x + Inches(2.8), Inches(3.9), Inches(1), Inches(0.2),
         font_size=8, color=GREEN)

# RTI recommendation
rounded_rect(sB, right_x + Inches(0.3), Inches(5.5), Inches(3.7), Inches(0.7),
             fill_color=RGBColor(0x14, 0x20, 0x0A), line_color=GREEN, line_w=1, radius=0.05)
add_text(sB, "\u2713  RECOMMENDATION: Move to RTI Tier 2", right_x + Inches(0.5), Inches(5.55), Inches(3.3), Inches(0.25),
         font_size=12, color=GREEN, bold=True)
add_text(sB, "Pattern: consistent decoding struggles at multisyllabic words", right_x + Inches(0.5), Inches(5.85), Inches(3.3), Inches(0.25),
         font_size=9, color=RGBColor(0x6B, 0x8A, 0x5A))

# --- PACER wordmark centered bottom ---
add_text(sB, "PACER", Inches(5.5), Inches(6.85), Inches(2.3), Inches(0.4),
         font_size=14, color=ACCENT1, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# CONCEPT C: "WHAT PACER HEARS" — Centered product, big passage hero
# Dark cinematic. One passage dominates. Words light up.
# ============================================================
sC = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sC, DARK_DEEP)

add_text(sC, "CONCEPT C", Inches(11.5), Inches(0.4), Inches(1.5), Inches(0.3),
         font_size=10, color=DIM, bold=True, alignment=PP_ALIGN.RIGHT)

# --- Top: PACER branding, huge ---
# Waveform bars - bigger
wave_bx = Inches(5.2)
wave_by = Inches(0.55)
for dx, h in [(0, 0.22), (0.14, 0.38), (0.28, 0.55), (0.42, 0.35), (0.56, 0.25)]:
    rounded_rect(sC, wave_bx + Inches(dx), wave_by + Inches((0.55-h)/2),
                 Inches(0.09), Inches(h), fill_color=ACCENT2, radius=0.5)

add_text(sC, "PACER", Inches(6.0), Inches(0.4), Inches(3), Inches(0.6),
         font_size=32, color=BRIGHT_WHITE, bold=True)

add_text(sC, "AI that hears how students struggle to read",
         Inches(3), Inches(1.15), Inches(7.3), Inches(0.4),
         font_size=18, color=MUTED, alignment=PP_ALIGN.CENTER)

# --- CENTER: Hero passage in a "window" ---
win_x = Inches(0.8)
win_y = Inches(1.9)
win_w = Inches(11.7)
win_h = Inches(4.0)

# Outer glow
rounded_rect(sC, win_x - Inches(0.03), win_y - Inches(0.03),
             win_w + Inches(0.06), win_h + Inches(0.06),
             fill_color=None, line_color=RGBColor(0x08, 0x50, 0x6B), line_w=0.5, radius=0.025)
# Window
rounded_rect(sC, win_x, win_y, win_w, win_h,
             fill_color=RGBColor(0x08, 0x0E, 0x1A), line_color=RGBColor(0x1E, 0x29, 0x3B), line_w=1, radius=0.025)

# Title bar
shape(sC, MSO_SHAPE.RECTANGLE, win_x, win_y, win_w, Inches(0.45),
      fill_color=RGBColor(0x0C, 0x14, 0x25))
for i, dc in enumerate([RED, RGBColor(0xEA, 0xB3, 0x08), GREEN]):
    shape(sC, MSO_SHAPE.OVAL, win_x + Inches(0.2 + i * 0.22), win_y + Inches(0.14),
          Inches(0.14), Inches(0.14), fill_color=dc)
add_text(sC, "pacer \u2014 live assessment \u2014 student: Alex M. \u2014 Grade 6",
         win_x + Inches(1.0), win_y + Inches(0.06), Inches(6), Inches(0.35),
         font_size=10, color=DIM, font_name='Consolas')

# --- The passage: BIG words, hero treatment ---
# Each word is a large, clearly readable element
passage_y_start = win_y + Inches(0.7)
# We'll use large font, well-spaced words
# 3 rows of passage text

# Row 1
r1_y = passage_y_start + Inches(0.1)
r1_words = [
    ("The", MUTED, 0.55, None), ("boy", MUTED, 0.6, None),
    ("walked", ORANGE, 1.1, "hesitation \u2022 720ms"),
    ("slowly", MUTED, 1.05, None), ("to", MUTED, 0.45, None),
    ("the", MUTED, 0.55, None),
    ("grocery", TEAL, 1.85, "STRUGGLE \u2022 3.2s \u2022 decoding failure"),
    ("store", MUTED, 0.85, None),
]

r2_y = passage_y_start + Inches(1.15)
r2_words = [
    ("and", MUTED, 0.6, None),
    ("bringed", ORANGE, 1.25, "substitution \u2022 bought \u2192 bringed"),
    ("some", MUTED, 0.85, None),
    ("\u23f8  4.1s", GRAY, 1.0, "long pause"),
    ("milk", PURPLE, 0.75, "self-corrected"),
    ("for", MUTED, 0.55, None), ("his", MUTED, 0.55, None),
    ("grandmother", RED, 2.1, "OMISSION \u2022 word skipped"),
]

r3_y = passage_y_start + Inches(2.2)
r3_words = [
    ("who", MUTED, 0.6, None), ("lived", MUTED, 0.85, None),
    ("around", ORANGE, 1.15, "substitution \u2022 across \u2192 around"),
    ("the", MUTED, 0.55, None), ("street.", MUTED, 1.0, None),
]

for ry, row in [(r1_y, r1_words), (r2_y, r2_words), (r3_y, r3_words)]:
    wx = win_x + Inches(0.5)
    for word, color, ww_in, annotation in row:
        ww = Inches(ww_in)

        # Annotation above (if present)
        if annotation:
            ann_color = color
            add_text(sC, annotation, wx, ry - Inches(0.28), ww + Inches(0.3), Inches(0.25),
                     font_size=8, color=ann_color, font_name='Consolas')

        # Word text
        if color == MUTED:
            # Dim correct words
            add_text(sC, word, wx, ry, ww, Inches(0.42),
                     font_size=26, color=RGBColor(0x50, 0x58, 0x68), font_name='Calibri')
        elif word.startswith("\u23f8"):
            # Pause indicator
            rounded_rect(sC, wx, ry + Inches(0.04), ww, Inches(0.35),
                         fill_color=RGBColor(0x15, 0x1A, 0x25), line_color=GRAY, line_w=0.75, radius=0.2)
            add_text(sC, word, wx + Inches(0.1), ry + Inches(0.02), ww - Inches(0.1), Inches(0.38),
                     font_size=16, color=GRAY, font_name='Consolas', alignment=PP_ALIGN.CENTER)
        else:
            # Colored word - bright, bold
            add_text(sC, word, wx, ry, ww, Inches(0.42),
                     font_size=26, color=color, bold=True, font_name='Calibri')
            # Underline
            if color == TEAL:
                # Dotted underline for struggle
                for di in range(int(ww_in / 0.12)):
                    if di % 2 == 0:
                        shape(sC, MSO_SHAPE.OVAL, wx + Inches(di * 0.12), ry + Inches(0.4),
                              Inches(0.05), Inches(0.05), fill_color=TEAL)
            elif color == RED:
                # Strikethrough for omission
                shape(sC, MSO_SHAPE.RECTANGLE, wx, ry + Inches(0.22), ww - Inches(0.1), Inches(0.025),
                      fill_color=RED)
            elif color != GRAY:
                # Solid underline
                shape(sC, MSO_SHAPE.RECTANGLE, wx, ry + Inches(0.4), ww - Inches(0.1), Inches(0.03),
                      fill_color=color)

        wx += ww + Inches(0.12)

# --- BOTTOM: Key stats + legend ---
bot_y = Inches(6.2)

# Legend
legend_items = [
    ("\u25CF Correct", RGBColor(0x50, 0x58, 0x68)),
    ("\u25CF Hesitation", ORANGE),
    ("\u25CF Substitution", ORANGE),
    ("\u25CF Struggle", TEAL),
    ("\u25CF Self-Correction", PURPLE),
    ("\u25CF Omission", RED),
    ("\u25CF Long Pause", GRAY),
]
for i, (label, color) in enumerate(legend_items):
    lx = Inches(0.8 + i * 1.7)
    add_text(sC, label, lx, bot_y, Inches(1.6), Inches(0.25),
             font_size=10, color=color, bold=True)

# Bottom tagline
add_text(sC, "13 types of reading struggle. Detected automatically. From a real book.",
         Inches(2.5), Inches(6.7), Inches(8.3), Inches(0.4),
         font_size=16, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

add_text(sC, "Middle School RTI Screening  \u2022  Zero Teacher Prep  \u2022  COPPA/FERPA Compliant",
         Inches(2.5), Inches(7.05), Inches(8.3), Inches(0.3),
         font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pitch-deck-hero-concepts.pptx')
prs.save(out)
print(f"Saved to {out}")
print("3 slides: Concept A (Before/After), Concept B (Book→Insight flow), Concept C (Cinematic passage hero)")
