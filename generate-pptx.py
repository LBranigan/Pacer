#!/usr/bin/env python3
"""Generate Pacer pitch deck PPTX file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG2 = RGBColor(0x1E, 0x29, 0x3B)
ACCENT1 = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x71, 0x16)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
TEAL = RGBColor(0x14, 0xB8, 0xA6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE,
             bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=16, color=MUTED, bullet_color=ACCENT1):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0
        pf = p._pPr
        if pf is None:
            from pptx.oxml.ns import qn
            from lxml import etree
            pf = etree.SubElement(p._p, qn('a:pPr'))
        pf.set('marL', str(Emu(Inches(0.3))))
        pf.set('indent', str(Emu(Inches(-0.25))))
    return txBox


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5), color=ACCENT1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, color=DARK_BG2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05
    return shape


def slide_number(slide, num, total=12):
    add_text(slide, f"{num} / {total}", Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
             font_size=10, color=DIM, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 - TITLE / HOOK
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s1, DARK_BG)

# Big faded "40%" background
add_text(s1, "40%", Inches(1), Inches(0.5), Inches(11), Inches(6.5),
         font_size=280, color=RGBColor(0x14, 0x1E, 0x33), bold=True, alignment=PP_ALIGN.CENTER)

add_text(s1, "The Forgotten 40%", Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
         font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(s1, "They don't qualify for Special Ed. They pass basic screenings.\nBut they devote 90% of their brainpower to decoding — leaving 0% for comprehension.",
         Inches(2.5), Inches(3.0), Inches(8), Inches(1.2),
         font_size=20, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(s1, "They've survived by memorizing words and guessing from context clues.\nA teacher might think, \"They're a little slow, but they're getting the words right.\"",
         Inches(2.5), Inches(4.4), Inches(8), Inches(1.0),
         font_size=16, color=DIM, italic=True, alignment=PP_ALIGN.CENTER)

# Pacer branding
add_text(s1, "PACER  —  AI-Powered Reading Struggle Detection", Inches(3), Inches(6.2), Inches(7), Inches(0.5),
         font_size=14, color=ACCENT1, alignment=PP_ALIGN.CENTER)

slide_number(s1, 1)

# ============================================================
# SLIDE 2 - THE PROBLEM
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, DARK_BG)

add_text(s2, "Reading Fluency Assessment Is Broken", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# Three cards
cards_data = [
    ("👩‍🏫", "Overwhelmed Teachers", "120+ students per teacher. Trained to teach content, not diagnose reading struggle. They need help, not more dashboards."),
    ("📊", "Blunt Metrics", "Current tools measure words-correct-per-minute — like diagnosing heart disease by taking temperature. Misses HOW students struggle."),
    ("⏳", "Closing Window", "Middle school is the last realistic intervention window. After 8th grade, outcomes calcify. No identification → no intervention."),
]

for i, (icon, title, desc) in enumerate(cards_data):
    x = Inches(0.8 + i * 4.1)
    add_card(s2, x, Inches(1.8), Inches(3.7), Inches(3.0))
    add_accent_bar(s2, x, Inches(1.8), Inches(3.7), Inches(0.04))
    add_text(s2, icon, x + Inches(0.3), Inches(2.1), Inches(1), Inches(0.5), font_size=28)
    add_text(s2, title, x + Inches(0.3), Inches(2.7), Inches(3.1), Inches(0.5), font_size=18, color=WHITE, bold=True)
    add_text(s2, desc, x + Inches(0.3), Inches(3.3), Inches(3.1), Inches(1.3), font_size=13, color=MUTED)

# Funnel
stages = [
    ("Elementary", "Some screening", GREEN),
    ("→", "", DIM),
    ("Middle School", "⚠ Gap — almost nothing", RED),
    ("→", "", DIM),
    ("High School", "Too late", DIM),
]
x_pos = Inches(2.0)
for label, sub, color in stages:
    if label == "→":
        add_text(s2, "→", x_pos, Inches(5.3), Inches(0.6), Inches(0.5), font_size=24, color=DIM, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.7)
    else:
        add_text(s2, label, x_pos, Inches(5.2), Inches(2.4), Inches(0.4), font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(s2, sub, x_pos, Inches(5.65), Inches(2.4), Inches(0.3), font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.5)

slide_number(s2, 2)

# ============================================================
# SLIDE 3 - WHY NOW
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, DARK_BG)

add_text(s3, "A LeNet Moment for Speech Recognition", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

bullets = [
    "ASR has reached a tipping point — but models are trained on adult, fluent speech",
    "Accurate ASR for disfluent populations (children, struggling readers) is the next frontier",
    "Edge compute makes COPPA/FERPA-compliant classroom deployment practical",
    "RTI is mandated but under-resourced — schools need tools, not more theory",
]
add_bullet_list(s3, bullets, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.5), font_size=16, color=MUTED)

# Bar chart - Adult vs Child
add_text(s3, "ASR Accuracy Gap", Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
         font_size=16, color=ACCENT2, bold=True)

# Adult bar
add_text(s3, "Adult Fluent Speech", Inches(7.2), Inches(2.5), Inches(3), Inches(0.3), font_size=12, color=MUTED)
add_text(s3, "~95%", Inches(11.5), Inches(2.5), Inches(1), Inches(0.3), font_size=12, color=ACCENT2)
bar_bg1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.9), Inches(5.3), Inches(0.45))
bar_bg1.fill.solid(); bar_bg1.fill.fore_color.rgb = DARK_BG2; bar_bg1.line.fill.background()
bar1 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.9), Inches(5.0), Inches(0.45))
bar1.fill.solid(); bar1.fill.fore_color.rgb = ACCENT1; bar1.line.fill.background()

# Child bar
add_text(s3, "Child Disfluent Speech", Inches(7.2), Inches(3.7), Inches(3), Inches(0.3), font_size=12, color=MUTED)
add_text(s3, "~35%", Inches(11.5), Inches(3.7), Inches(1), Inches(0.3), font_size=12, color=RED)
bar_bg2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.1), Inches(5.3), Inches(0.45))
bar_bg2.fill.solid(); bar_bg2.fill.fore_color.rgb = DARK_BG2; bar_bg2.line.fill.background()
bar2 = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(4.1), Inches(1.9), Inches(0.45))
bar2.fill.solid(); bar2.fill.fore_color.rgb = RED; bar2.line.fill.background()

# Gap callout
gap_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(4.9), Inches(4.5), Inches(0.6))
gap_box.fill.solid(); gap_box.fill.fore_color.rgb = RGBColor(0x0B, 0x1D, 0x35)
gap_box.line.color.rgb = ACCENT1; gap_box.line.width = Pt(1)
gap_box.line.dash_style = 4  # dash
add_text(s3, "↕  This gap is the entire opportunity", Inches(7.8), Inches(4.95), Inches(4), Inches(0.45),
         font_size=15, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

slide_number(s3, 3)

# ============================================================
# SLIDE 4 - SOLUTION
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, DARK_BG)

add_text(s4, "A Struggle Detector, Not a Score Generator", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# Flow steps
flow = [("📖", "Real Books", "Student reads a\nphysical book"),
        ("🎙️", "Pacer Listens", "Classroom device\ncaptures audio"),
        ("🧠", "AI Analyzes", "Multi-engine struggle\nclassification"),
        ("📋", "Teacher Sees", "Rich data, zero\nprep required")]
for i, (icon, title, desc) in enumerate(flow):
    x = Inches(0.8 + i * 3.2)
    add_card(s4, x, Inches(1.7), Inches(2.6), Inches(2.4))
    add_text(s4, icon, x + Inches(0.2), Inches(1.9), Inches(0.6), Inches(0.5), font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(s4, title, x + Inches(0.2), Inches(2.5), Inches(2.2), Inches(0.4), font_size=16, color=WHITE, bold=True)
    add_text(s4, desc, x + Inches(0.2), Inches(3.0), Inches(2.2), Inches(0.8), font_size=12, color=MUTED)
    if i < 3:
        add_text(s4, "→", x + Inches(2.6), Inches(2.5), Inches(0.5), Inches(0.5), font_size=24, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Quote box
add_accent_bar(s4, Inches(0.8), Inches(4.6), Inches(0.06), Inches(0.6))
add_text(s4, '"A blood pressure cuff for reading — installed in the classroom, reliable, unobtrusive, always on."',
         Inches(1.1), Inches(4.55), Inches(10), Inches(0.7), font_size=18, color=WHITE, italic=True)

# Feature chips
features = ["Hesitations", "Substitutions", "Repetitions", "Omissions", "Self-corrections", "Prosody", "Longitudinal trends"]
for i, feat in enumerate(features):
    x = Inches(0.8 + i * 1.7)
    add_text(s4, f"●  {feat}", x, Inches(5.6), Inches(1.6), Inches(0.4), font_size=11, color=MUTED)

slide_number(s4, 4)

# ============================================================
# SLIDE 5 - HOW IT WORKS
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s5, DARK_BG)

add_text(s5, "13-Point Miscue Classification", Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
         font_size=42, color=WHITE, bold=True)

# Pipeline
pipe_steps = ["🎤 Audio", "Multi-Engine ASR\n(3 engines)", "Forced\nAlignment", "Disfluency\nDetection", "Struggle\nClassification", "📊 Dashboard"]
for i, step in enumerate(pipe_steps):
    x = Inches(0.5 + i * 2.1)
    c = ACCENT1 if i in [1, 3, 5] else DARK_BG2
    box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(1.8), Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = c
    box.line.color.rgb = RGBColor(0x33, 0x41, 0x55) if i not in [1,3,5] else ACCENT2
    box.line.width = Pt(1)
    add_text(s5, step, x + Inches(0.05), Inches(1.45), Inches(1.7), Inches(0.8),
             font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s5, "→", x + Inches(1.8), Inches(1.55), Inches(0.3), Inches(0.5), font_size=16, color=ACCENT1)

# Miscue categories
add_text(s5, "COUNTED AS ERRORS", Inches(0.8), Inches(2.8), Inches(4), Inches(0.3),
         font_size=10, color=DIM, bold=True)
errors = ["Omission", "Substitution", "Struggle", "Morphological", "Long Pause"]
error_colors = [RED, ORANGE, TEAL, ORANGE, DIM]
for i, (tag, col) in enumerate(zip(errors, error_colors)):
    x = Inches(0.8 + i * 2.0)
    add_text(s5, f"■ {tag}", x, Inches(3.2), Inches(1.8), Inches(0.3), font_size=13, color=col, bold=True)

add_text(s5, "DIAGNOSTIC (SHOWN, NOT SCORED)", Inches(0.8), Inches(3.8), Inches(6), Inches(0.3),
         font_size=10, color=DIM, bold=True)
diags = ["Insertion", "Hesitation", "Self-Correction", "Fragments", "Repetitions", "Fillers"]
diag_colors = [RGBColor(0x3B, 0x82, 0xF6), ORANGE, PURPLE, RGBColor(0xEC, 0x48, 0x99), RGBColor(0xEC, 0x48, 0x99), PURPLE]
for i, (tag, col) in enumerate(zip(diags, diag_colors)):
    x = Inches(0.8 + i * 2.0)
    add_text(s5, f"■ {tag}", x, Inches(4.2), Inches(1.8), Inches(0.3), font_size=13, color=col, bold=True)

add_text(s5, "FORGIVENESS RULES", Inches(0.8), Inches(4.8), Inches(4), Inches(0.3),
         font_size=10, color=DIM, bold=True)
add_text(s5, "■ Proper Noun Recognition", Inches(0.8), Inches(5.2), Inches(3), Inches(0.3),
         font_size=13, color=GREEN, bold=True)

# Engine badge
add_text(s5, "⚡  3-Engine Cross-Validation:  Reverb  •  Google STT  •  Deepgram Nova-3",
         Inches(0.8), Inches(5.9), Inches(8), Inches(0.4), font_size=14, color=ACCENT2, bold=True)

slide_number(s5, 5)

# ============================================================
# SLIDE 6 - MARKET
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s6, DARK_BG)

add_text(s6, "$2.4B Market — Massive Gap in the Middle", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# Concentric circles (using shapes)
cx, cy = Inches(3.2), Inches(3.8)
# Outer
outer = s6.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(2.5), cy - Inches(2.5), Inches(5), Inches(5))
outer.fill.solid(); outer.fill.fore_color.rgb = RGBColor(0x12, 0x1B, 0x30)
outer.line.color.rgb = RGBColor(0x33, 0x41, 0x55); outer.line.width = Pt(1)
add_text(s6, "Expansion", cx - Inches(2.3), cy - Inches(2.3), Inches(2), Inches(0.3), font_size=10, color=DIM)
add_text(s6, "SLP • Private • Homeschool", cx - Inches(1.5), cy + Inches(1.8), Inches(3), Inches(0.3), font_size=9, color=DIM, alignment=PP_ALIGN.CENTER)
# Mid
mid = s6.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(1.7), cy - Inches(1.7), Inches(3.4), Inches(3.4))
mid.fill.solid(); mid.fill.fore_color.rgb = RGBColor(0x0C, 0x23, 0x3D)
mid.line.color.rgb = RGBColor(0x0E, 0xA5, 0xE9); mid.line.width = Pt(1)
add_text(s6, "Adjacent", cx - Inches(1.5), cy - Inches(1.5), Inches(1.5), Inches(0.3), font_size=10, color=ACCENT2)
add_text(s6, "K-5 • High School", cx - Inches(1), cy + Inches(0.9), Inches(2), Inches(0.3), font_size=9, color=ACCENT2, alignment=PP_ALIGN.CENTER)
# Core
core = s6.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.9), cy - Inches(0.9), Inches(1.8), Inches(1.8))
core.fill.solid(); core.fill.fore_color.rgb = RGBColor(0x08, 0x3D, 0x51)
core.line.color.rgb = ACCENT1; core.line.width = Pt(2)
add_text(s6, "Core:\nMiddle School\nRTI", cx - Inches(0.7), cy - Inches(0.5), Inches(1.4), Inches(1), font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Market points on the right
points = [
    "🎯  Middle school is drastically underserved — fewer competitors, desperate need",
    "🏫  B2B entry via RTI/MTSS mandates — existing budget line items",
    "👨‍👩‍👦  B2C expansion: SLPs, private schools, homeschool families",
    "🗄️  Data moat: every session = proprietary disfluent children's speech data — scarcest dataset in edtech ASR",
]
for i, point in enumerate(points):
    y = Inches(1.8 + i * 1.2)
    add_text(s6, point, Inches(6.5), y, Inches(6), Inches(1.0), font_size=15, color=MUTED)

slide_number(s6, 6)

# ============================================================
# SLIDE 7 - COMPETITION
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s7, DARK_BG)

add_text(s7, "Positioned Where No One Else Is Playing", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# 2x2 Matrix
mx, my = Inches(2.5), Inches(1.8)
mw, mh = Inches(4.0), Inches(2.4)
gap = Inches(0.06)

# Top-left: Screen + Score (MAP, iStation)
c1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, my, mw, mh)
c1.fill.solid(); c1.fill.fore_color.rgb = DARK_BG2; c1.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
add_text(s7, "NWEA MAP\nLegacy • Standardized scores", mx + Inches(0.3), my + Inches(0.5), Inches(3), Inches(0.8), font_size=14, color=MUTED)
add_text(s7, "iStation\nLegacy • WCPM only", mx + Inches(0.3), my + Inches(1.4), Inches(3), Inches(0.8), font_size=14, color=MUTED)

# Top-right: Screen + Struggle (empty)
c2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx + mw + gap, my, mw, mh)
c2.fill.solid(); c2.fill.fore_color.rgb = DARK_BG2; c2.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
add_text(s7, "—", mx + mw + gap + Inches(1.5), my + Inches(0.9), Inches(1), Inches(0.5), font_size=18, color=DIM, alignment=PP_ALIGN.CENTER)

# Bottom-left: Books + Score (Amira, Ello)
c3 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx, my + mh + gap, mw, mh)
c3.fill.solid(); c3.fill.fore_color.rgb = DARK_BG2; c3.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
add_text(s7, "Amira Learning\n~$40M raised • AI tutoring", mx + Inches(0.3), my + mh + gap + Inches(0.5), Inches(3), Inches(0.8), font_size=14, color=MUTED)
add_text(s7, "Ello\nReed Hastings, YC • K-2 books", mx + Inches(0.3), my + mh + gap + Inches(1.4), Inches(3), Inches(0.8), font_size=14, color=MUTED)

# Bottom-right: Books + Struggle = PACER (highlighted!)
c4 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx + mw + gap, my + mh + gap, mw, mh)
c4.fill.solid(); c4.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x3F)
c4.line.color.rgb = ACCENT2; c4.line.width = Pt(2)
add_text(s7, "PACER", mx + mw + gap + Inches(0.8), my + mh + gap + Inches(0.6), Inches(2.4), Inches(0.6),
         font_size=32, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s7, "Alone in this quadrant", mx + mw + gap + Inches(0.8), my + mh + gap + Inches(1.3), Inches(2.4), Inches(0.4),
         font_size=13, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Axis labels
add_text(s7, "Score-level  →  Struggle-level granularity", Inches(4), Inches(6.8), Inches(5), Inches(0.3),
         font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(s7, "Screen-based", mx - Inches(1.6), my + Inches(0.8), Inches(1.5), Inches(0.4),
         font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)
add_text(s7, "Physical books", mx - Inches(1.6), my + mh + gap + Inches(0.8), Inches(1.5), Inches(0.4),
         font_size=11, color=DIM, alignment=PP_ALIGN.CENTER)

slide_number(s7, 7)

# ============================================================
# SLIDE 8 - TRACTION
# ============================================================
s8 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s8, DARK_BG)

add_text(s8, "Already in Classrooms", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# Timeline
tl_items = [
    ("✓", "Built MVP", "Full AI pipeline live", False),
    ("◉", "Morningside Pilot", "Active — real students", True),
    ("3-5", "School Expansion", "Precision teaching network", False),
    ("⟶", "District Partners", "Paid pilots", False),
]
for i, (icon, title, desc, active) in enumerate(tl_items):
    x = Inches(1.0 + i * 3.0)
    dot_color = ACCENT1 if active else DARK_BG2
    dot = s8.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), Inches(1.8), Inches(0.5), Inches(0.5))
    dot.fill.solid(); dot.fill.fore_color.rgb = dot_color
    dot.line.color.rgb = ACCENT1 if active else RGBColor(0x33, 0x41, 0x55)
    dot.line.width = Pt(2)
    add_text(s8, icon, x + Inches(0.82), Inches(1.82), Inches(0.46), Inches(0.46),
             font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s8, title, x, Inches(2.5), Inches(2.2), Inches(0.4), font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s8, desc, x, Inches(2.95), Inches(2.2), Inches(0.3), font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)
    if i < 3:
        line = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(1.3), Inches(2.03), Inches(2.0), Inches(0.03))
        line.fill.solid(); line.fill.fore_color.rgb = DARK_BG2; line.line.fill.background()

# Traction cards
traction = [
    ("🏫", "Morningside Academy — live pilot, real students, real data, zero bureaucratic friction"),
    ("🔗", "Andrew Kieta — precision teaching network opens doors to data-driven schools"),
    ("⚙️", "Product functional today — Word Speed Map, disfluency detection, multi-miscue engine all live"),
    ("🔒", "COPPA/FERPA compliant architecture from day one — edge compute, no cloud-stored student audio"),
]
for i, (icon, text) in enumerate(traction):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(3.7 + row * 1.4)
    add_card(s8, x, y, Inches(5.8), Inches(1.1))
    add_text(s8, icon, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.5), font_size=20)
    add_text(s8, text, x + Inches(0.8), y + Inches(0.2), Inches(4.7), Inches(0.7), font_size=13, color=MUTED)

slide_number(s8, 8)

# ============================================================
# SLIDE 9 - TEAM
# ============================================================
s9 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s9, DARK_BG)

add_text(s9, "Built for This Problem", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

# Emma card
add_card(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(3.5))
add_accent_bar(s9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.05))
add_text(s9, "Emma", Inches(1.2), Inches(1.9), Inches(3), Inches(0.5), font_size=24, color=WHITE, bold=True)
emma_creds = [
    "PhD from Vanderbilt — #1 Special Education program",
    "Mentored by Doug Fuchs (creator of PALs system)",
    "Dual expertise: Linguistics + Special Education",
    "Decade of classroom teaching experience",
    "Currently runs the Wing Institute",
]
for i, cred in enumerate(emma_creds):
    add_text(s9, f"•  {cred}", Inches(1.2), Inches(2.6 + i * 0.42), Inches(4.8), Inches(0.4), font_size=13, color=MUTED)

# Founder card
add_card(s9, Inches(6.8), Inches(1.6), Inches(5.6), Inches(3.5))
add_accent_bar(s9, Inches(6.8), Inches(1.6), Inches(5.6), Inches(0.05))
add_text(s9, "Founder", Inches(7.2), Inches(1.9), Inches(3), Inches(0.5), font_size=24, color=WHITE, bold=True)
founder_creds = [
    "10+ years as entrepreneur",
    "Built Pacer's entire AI pipeline solo",
    "All-in: 12+ hour days, 7 days/week",
    "Deep technical: AI/ML, edge compute, audio",
    "Former Amazon seller — understands PMF",
]
for i, cred in enumerate(founder_creds):
    add_text(s9, f"•  {cred}", Inches(7.2), Inches(2.6 + i * 0.42), Inches(4.8), Inches(0.4), font_size=13, color=MUTED)

# Network bar
add_card(s9, Inches(0.8), Inches(5.5), Inches(5.6), Inches(1.2))
add_text(s9, "🏫  Morningside Academy", Inches(1.2), Inches(5.7), Inches(4), Inches(0.3), font_size=14, color=WHITE, bold=True)
add_text(s9, "Free classroom access • Pilot partner • Network gateway", Inches(1.2), Inches(6.1), Inches(4), Inches(0.3), font_size=11, color=MUTED)

add_card(s9, Inches(6.8), Inches(5.5), Inches(5.6), Inches(1.2))
add_text(s9, "💻  Nvidia Connection", Inches(7.2), Inches(5.7), Inches(4), Inches(0.3), font_size=14, color=WHITE, bold=True)
add_text(s9, "Connor — Distinguished Engineer • GPU/inference expertise", Inches(7.2), Inches(6.1), Inches(4), Inches(0.3), font_size=11, color=MUTED)

slide_number(s9, 9)

# ============================================================
# SLIDE 10 - BUSINESS MODEL
# ============================================================
s10 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s10, DARK_BG)

add_text(s10, "Land with Screening, Expand with Data", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)

stages = [
    ("01 / LAND", "RTI Screening", "Per-school SaaS subscription.\nAligns with existing budget\nline items and mandates.", True),
    ("02 / EXPAND", "District Analytics", "Progress monitoring, longitudinal\ndashboards, tier 2/3 intervention\ntracking.", False),
    ("03 / DEFEND", "Data Moat", "Proprietary disfluent speech data.\nLicense fine-tuned model weights\nto other platforms.", False),
]
for i, (num, title, desc, primary) in enumerate(stages):
    x = Inches(0.8 + i * 4.1)
    border_col = ACCENT2 if primary else RGBColor(0x33, 0x41, 0x55)
    card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.7), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x3F) if primary else DARK_BG2
    card.line.color.rgb = border_col; card.line.width = Pt(1)
    add_text(s10, num, x + Inches(0.3), Inches(1.9), Inches(3), Inches(0.3), font_size=11, color=ACCENT1, bold=True)
    add_text(s10, title, x + Inches(0.3), Inches(2.3), Inches(3), Inches(0.4), font_size=22, color=WHITE, bold=True)
    add_text(s10, desc, x + Inches(0.3), Inches(2.9), Inches(3), Inches(1.2), font_size=13, color=MUTED)

# Soapbox callout
callout = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.2))
callout.fill.solid(); callout.fill.fore_color.rgb = RGBColor(0x1A, 0x12, 0x12)
callout.line.color.rgb = RGBColor(0x4B, 0x1D, 0x1D); callout.line.width = Pt(1)
add_text(s10, "💡", Inches(1.1), Inches(5.15), Inches(0.5), Inches(0.5), font_size=22)
add_text(s10, "Soapbox Labs was acquired by Curriculum Associates for an estimated $100M+ — for the children's speech data. First thing they did: closed the API to competitors. Our data will be richer because we capture disfluency, not just correct speech.",
         Inches(1.7), Inches(5.1), Inches(10.5), Inches(0.9), font_size=14, color=MUTED)

slide_number(s10, 10)

# ============================================================
# SLIDE 11 - THE ASK
# ============================================================
s11 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s11, DARK_BG)

add_text(s11, "Accelerate Pacer into 50 Classrooms", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         font_size=42, color=WHITE, bold=True)
add_text(s11, "Seeking Launch.co Partnership", Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
         font_size=20, color=MUTED, alignment=PP_ALIGN.CENTER)

# Three ask cards
asks = [
    ("🧑‍💻", "Hire", "First ML engineer to accelerate\nmodel development"),
    ("🖥️", "Deploy", "Edge hardware in classrooms\nfor on-device inference"),
    ("📝", "Validate", "3 paid district pilots to\nprove sales motion"),
]
for i, (icon, title, desc) in enumerate(asks):
    x = Inches(0.8 + i * 4.1)
    add_card(s11, x, Inches(2.2), Inches(3.7), Inches(2.2))
    add_text(s11, icon, x + Inches(1.3), Inches(2.4), Inches(1), Inches(0.5), font_size=28, alignment=PP_ALIGN.CENTER)
    add_text(s11, title, x + Inches(0.3), Inches(3.0), Inches(3.1), Inches(0.4), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s11, desc, x + Inches(0.3), Inches(3.5), Inches(3.1), Inches(0.8), font_size=13, color=MUTED, alignment=PP_ALIGN.CENTER)

# Milestones bar
add_card(s11, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5))
milestones = [("50", "Classrooms"), ("3", "District LOIs"), ("1", "Validation Study"), ("12 mo", "Timeline")]
for i, (num, label) in enumerate(milestones):
    x = Inches(1.8 + i * 2.8)
    add_text(s11, num, x, Inches(5.2), Inches(1.8), Inches(0.6), font_size=36 if len(num) <= 2 else 20, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(s11, label, x, Inches(5.85), Inches(1.8), Inches(0.3), font_size=11, color=MUTED, alignment=PP_ALIGN.CENTER)

slide_number(s11, 11)

# ============================================================
# SLIDE 12 - CLOSING
# ============================================================
s12 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s12, DARK_BG)

add_text(s12, "Every Struggling Reader Found.", Inches(1), Inches(2.2), Inches(11), Inches(1),
         font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(s12, "None Left Behind.", Inches(1), Inches(3.4), Inches(11), Inches(0.8),
         font_size=36, color=MUTED, alignment=PP_ALIGN.CENTER)

add_text(s12, "PACER  —  AI-Powered Reading Struggle Detection", Inches(3), Inches(5.2), Inches(7), Inches(0.4),
         font_size=14, color=ACCENT1, alignment=PP_ALIGN.CENTER)

# Contact placeholder
add_text(s12, "yourname@email.com  •  yourwebsite.com", Inches(3), Inches(5.8), Inches(7), Inches(0.4),
         font_size=13, color=DIM, alignment=PP_ALIGN.CENTER)

slide_number(s12, 12)

# ============================================================
# SAVE
# ============================================================
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'pitch-deck.pptx')
prs.save(out)
print(f"Saved to {out}")
